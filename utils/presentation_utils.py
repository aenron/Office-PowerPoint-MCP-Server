"""
Presentation management utilities for PowerPoint MCP Server.
Functions for creating, opening, saving, and managing presentations.
"""
from pptx import Presentation
from datetime import date, datetime
from typing import Any, Dict, List, Optional, Tuple
import os


def create_presentation() -> Presentation:
    """
    Create a new PowerPoint presentation.

    Returns:
        A new Presentation object
    """
    return Presentation()


def open_presentation(file_path: str) -> Presentation:
    """
    Open an existing PowerPoint presentation.

    Args:
        file_path: Path to the PowerPoint file

    Returns:
        A Presentation object
    """
    return Presentation(file_path)


def create_presentation_from_template(template_path: str) -> Presentation:
    """
    Create a new PowerPoint presentation from a template file.

    Args:
        template_path: Path to the template .pptx file

    Returns:
        A new Presentation object based on the template

    Raises:
        FileNotFoundError: If the template file doesn't exist
        Exception: If the template file is corrupted or invalid
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")

    if not template_path.lower().endswith(('.pptx', '.potx')):
        raise ValueError("Template file must be a .pptx or .potx file")

    try:
        # Load the template file as a presentation
        presentation = Presentation(template_path)
        return presentation
    except Exception as e:
        raise Exception(
            f"Failed to load template file '{template_path}': {str(e)}")


def save_presentation(presentation: Presentation, file_path: str) -> str:
    """
    Save a PowerPoint presentation to a file.

    Args:
        presentation: The Presentation object
        file_path: Path where the file should be saved

    Returns:
        The file path where the presentation was saved
    """
    presentation.save(file_path)
    return file_path


def cleanup_stale_generated_files(
    directory: str,
    keep_date: Optional[date] = None,
    suffixes: Tuple[str, ...] = (".pptx",),
) -> Dict[str, Any]:
    """
    Delete generated presentation files whose modified date is not keep_date.

    Args:
        directory: Directory containing generated files.
        keep_date: Local date to retain. Defaults to today.
        suffixes: File suffixes considered generated presentation files.

    Returns:
        Cleanup summary. Delete failures are reported but not raised.
    """
    retention_date = keep_date or date.today()
    summary: Dict[str, Any] = {
        "directory": os.path.abspath(directory) if directory else "",
        "retention_date": retention_date.isoformat(),
        "deleted_files": [],
        "failed_files": [],
    }
    if not directory or not os.path.isdir(directory):
        return summary

    normalized_suffixes = tuple(suffix.lower() for suffix in suffixes)
    for entry in os.scandir(directory):
        if not entry.is_file():
            continue
        if normalized_suffixes and not entry.name.lower().endswith(normalized_suffixes):
            continue

        try:
            modified_date = datetime.fromtimestamp(entry.stat().st_mtime).date()
        except OSError as exc:
            summary["failed_files"].append({
                "file_path": entry.path,
                "error": str(exc),
            })
            continue

        if modified_date == retention_date:
            continue

        try:
            os.remove(entry.path)
            summary["deleted_files"].append(entry.path)
        except OSError as exc:
            summary["failed_files"].append({
                "file_path": entry.path,
                "error": str(exc),
            })

    return summary


def is_file_modified_on_date(file_path: str, keep_date: Optional[date] = None) -> bool:
    """Return whether file_path exists and was modified on keep_date."""
    if not file_path or not os.path.isfile(file_path):
        return False
    retention_date = keep_date or date.today()
    try:
        return datetime.fromtimestamp(os.path.getmtime(file_path)).date() == retention_date
    except OSError:
        return False


def get_template_info(template_path: str) -> Dict:
    """
    Get information about a template file.

    Args:
        template_path: Path to the template .pptx file

    Returns:
        Dictionary containing template information
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")

    try:
        presentation = Presentation(template_path)

        # Get slide layouts
        layouts = get_slide_layouts(presentation)

        # Get core properties
        core_props = get_core_properties(presentation)

        # Get slide count
        slide_count = len(presentation.slides)

        # Get file size
        file_size = os.path.getsize(template_path)

        return {
            "template_path": template_path,
            "file_size_bytes": file_size,
            "slide_count": slide_count,
            "layout_count": len(layouts),
            "slide_layouts": layouts,
            "core_properties": core_props
        }
    except Exception as e:
        raise Exception(
            f"Failed to read template info from '{template_path}': {str(e)}")


def get_presentation_info(presentation: Presentation) -> Dict:
    """
    Get information about a presentation.

    Args:
        presentation: The Presentation object

    Returns:
        Dictionary containing presentation information
    """
    try:
        # Get slide layouts
        layouts = get_slide_layouts(presentation)

        # Get core properties
        core_props = get_core_properties(presentation)

        # Get slide count
        slide_count = len(presentation.slides)

        return {
            "slide_count": slide_count,
            "layout_count": len(layouts),
            "slide_layouts": layouts,
            "core_properties": core_props,
            "slide_width": presentation.slide_width,
            "slide_height": presentation.slide_height
        }
    except Exception as e:
        raise Exception(f"Failed to get presentation info: {str(e)}")


def get_slide_layouts(presentation: Presentation) -> List[Dict]:
    """
    Get all available slide layouts in the presentation.

    Args:
        presentation: The Presentation object

    Returns:
        A list of dictionaries with layout information
    """
    layouts = []
    for i, layout in enumerate(presentation.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholder_count": len(layout.placeholders)
        }
        layouts.append(layout_info)
    return layouts


def set_core_properties(presentation: Presentation, title: str = None, subject: str = None,
                        author: str = None, keywords: str = None, comments: str = None) -> None:
    """
    Set core document properties.

    Args:
        presentation: The Presentation object
        title: Document title
        subject: Document subject
        author: Document author
        keywords: Document keywords
        comments: Document comments
    """
    core_props = presentation.core_properties

    if title is not None:
        core_props.title = title
    if subject is not None:
        core_props.subject = subject
    if author is not None:
        core_props.author = author
    if keywords is not None:
        core_props.keywords = keywords
    if comments is not None:
        core_props.comments = comments


def get_core_properties(presentation: Presentation) -> Dict:
    """
    Get core document properties.

    Args:
        presentation: The Presentation object

    Returns:
        Dictionary containing core properties
    """
    core_props = presentation.core_properties

    return {
        "title": core_props.title,
        "subject": core_props.subject,
        "author": core_props.author,
        "keywords": core_props.keywords,
        "comments": core_props.comments,
        "created": core_props.created.isoformat() if core_props.created else None,
        "last_modified_by": core_props.last_modified_by,
        "modified": core_props.modified.isoformat() if core_props.modified else None
    }
