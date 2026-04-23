"""
Content management tools for PowerPoint MCP Server.
Handles slides, text, images, and content manipulation.
"""
from typing import Dict, List, Any, Union
from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations
import utils as ppt_utils
import tempfile
import base64
import os


def register_content_tools(app: FastMCP, presentations: Dict, get_current_presentation_id, validate_parameters, is_positive, is_non_negative, is_in_range, is_valid_rgb):
    """Register content management tools with the FastMCP app"""

    @app.tool(
        annotations=ToolAnnotations(
            title="Add Slide",
        ),
    )
    def add_slide(
        presentation_id: str,
        layout_index: int = 1,
        title: str = "",
        # "solid", "gradient", "professional_gradient"
        background_type: str = "",
        # For gradient: [[start_rgb], [end_rgb]]
        background_colors: List[List[int]] = [],
        gradient_direction: str = "horizontal",
        color_scheme: str = "modern_blue",
    ) -> Dict:
        """Add a new slide to the presentation with optional background styling."""
        if presentation_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = presentations[presentation_id]

        # Validate layout index
        if layout_index < 0 or layout_index >= len(pres.slide_layouts):
            return {
                "error": f"Invalid layout index: {layout_index}. Available layouts: 0-{len(pres.slide_layouts) - 1}"
            }

        try:
            # Add the slide
            slide, layout = ppt_utils.add_slide(pres, layout_index)
            slide_index = len(pres.slides) - 1

            # Set title if provided
            if title:
                ppt_utils.set_title(slide, title)

            # Apply background if specified
            if background_type == "gradient" and background_colors and len(background_colors) >= 2:
                ppt_utils.set_slide_gradient_background(
                    slide, background_colors[0], background_colors[1], gradient_direction
                )
            elif background_type == "professional_gradient":
                ppt_utils.create_professional_gradient_background(
                    slide, color_scheme, "subtle", gradient_direction
                )

            return {
                "message": f"Added slide {slide_index} with layout {layout_index}",
                "slide_index": slide_index,
                "layout_name": layout.name if hasattr(layout, 'name') else f"Layout {layout_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to add slide: {str(e)}"
            }

    @app.tool(
        annotations=ToolAnnotations(
            title="Get Slide Info",
            readOnlyHint=True,
        ),
    )
    def get_slide_info(presentation_id: str, slide_index: int) -> Dict:
        """Get information about a specific slide."""
        if presentation_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = presentations[presentation_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            return ppt_utils.get_slide_info(slide, slide_index)
        except Exception as e:
            return {
                "error": f"Failed to get slide info: {str(e)}"
            }

    @app.tool(
        annotations=ToolAnnotations(
            title="Extract Slide Text",
            readOnlyHint=True,
        ),
    )
    def extract_slide_text(presentation_id: str, slide_index: int) -> Dict:
        """Extract all text content from a specific slide."""
        if presentation_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = presentations[presentation_id]

        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        slide = pres.slides[slide_index]

        try:
            result = ppt_utils.extract_slide_text_content(slide)
            result["slide_index"] = slide_index
            return result
        except Exception as e:
            return {
                "error": f"Failed to extract slide text: {str(e)}"
            }

    @app.tool(
        annotations=ToolAnnotations(
            title="Extract Presentation Text",
            readOnlyHint=True,
        ),
    )
    def extract_presentation_text(presentation_id: str, include_slide_info: bool = True) -> Dict:
        """Extract all text content from all slides in the presentation."""
        if presentation_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }

        pres = presentations[presentation_id]

        try:
            slides_text = []
            total_text_shapes = 0
            slides_with_tables = 0
            slides_with_titles = 0
            all_presentation_text = []

            for slide_index, slide in enumerate(pres.slides):
                slide_text_result = ppt_utils.extract_slide_text_content(slide)

                if slide_text_result["success"]:
                    slide_data = {
                        "slide_index": slide_index,
                        "text_content": slide_text_result["text_content"]
                    }

                    if include_slide_info:
                        # Add basic slide info
                        slide_data["layout_name"] = slide.slide_layout.name
                        slide_data["total_text_shapes"] = slide_text_result["total_text_shapes"]
                        slide_data["has_title"] = slide_text_result["has_title"]
                        slide_data["has_tables"] = slide_text_result["has_tables"]

                    slides_text.append(slide_data)

                    # Accumulate statistics
                    total_text_shapes += slide_text_result["total_text_shapes"]
                    if slide_text_result["has_tables"]:
                        slides_with_tables += 1
                    if slide_text_result["has_title"]:
                        slides_with_titles += 1

                    # Collect all text for combined output
                    if slide_text_result["text_content"]["all_text_combined"]:
                        all_presentation_text.append(
                            f"=== SLIDE {slide_index + 1} ===")
                        all_presentation_text.append(
                            slide_text_result["text_content"]["all_text_combined"])
                        all_presentation_text.append(
                            "")  # Empty line separator
                else:
                    slides_text.append({
                        "slide_index": slide_index,
                        "error": slide_text_result.get("error", "Unknown error"),
                        "text_content": None
                    })

            return {
                "success": True,
                "presentation_id": presentation_id,
                "total_slides": len(pres.slides),
                "slides_with_text": len([s for s in slides_text if s.get("text_content") is not None]),
                "total_text_shapes": total_text_shapes,
                "slides_with_titles": slides_with_titles,
                "slides_with_tables": slides_with_tables,
                "slides_text": slides_text,
                "all_presentation_text_combined": "\n".join(all_presentation_text)
            }

        except Exception as e:
            return {
                "error": f"Failed to extract presentation text: {str(e)}"
            }

    @app.tool(
        annotations=ToolAnnotations(
            title="Populate Placeholder",
        ),
    )
    def populate_placeholder(
        presentation_id: str,
        slide_index: int,
        placeholder_idx: int,
        text: str,
    ) -> Dict:
        """Populate a placeholder with text."""
        slide, error = get_required_slide(presentation_id, slide_index)
        if error:
            return error

        try:
            ppt_utils.populate_placeholder(slide, placeholder_idx, text)
            return {
                "message": f"Populated placeholder {placeholder_idx} on slide {slide_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to populate placeholder: {str(e)}"
            }

    @app.tool(
        annotations=ToolAnnotations(
            title="Add Bullet Points",
        ),
    )
    def add_bullet_points(
        presentation_id: str,
        slide_index: int,
        placeholder_idx: int,
        bullet_points: List[str],
    ) -> Dict:
        """Add bullet points to a placeholder."""
        slide, error = get_required_slide(presentation_id, slide_index)
        if error:
            return error

        try:
            placeholder = slide.placeholders[placeholder_idx]
            ppt_utils.add_bullet_points(placeholder, bullet_points)
            return {
                "message": f"Added {len(bullet_points)} bullet points to placeholder {placeholder_idx} on slide {slide_index}"
            }
        except Exception as e:
            return {
                "error": f"Failed to add bullet points: {str(e)}"
            }

    def get_required_presentation(presentation_id: str):
        if presentation_id not in presentations:
            return None, {
                "error": "The specified presentation ID is invalid"
            }
        return presentations[presentation_id], None

    def get_required_slide(presentation_id: str, slide_index: int):
        pres, error = get_required_presentation(presentation_id)
        if error:
            return None, error

        if slide_index < 0 or slide_index >= len(pres.slides):
            return None, {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }

        return pres.slides[slide_index], None

    def get_required_shape(presentation_id: str, slide_index: int, shape_index: int):
        slide, error = get_required_slide(presentation_id, slide_index)
        if error:
            return None, error

        if shape_index < 0 or shape_index >= len(slide.shapes):
            return None, {
                "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
            }

        return slide.shapes[shape_index], None

    @app.tool(
        annotations=ToolAnnotations(
            title="Manage Text Box",
        ),
    )
    def manage_text_box(
        presentation_id: str,
        operation: str,
        slide_index: int,
        shape_index: int,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_size: int,
        font_name: str,
        bold: bool,
        italic: bool,
        underline: bool,
        color: List[int],
        bg_color: List[int],
        alignment: str,
        vertical_alignment: str,
        text_runs: List[Dict],
        min_font_size: int,
        max_font_size: int,
    ) -> Dict:
        """Manage text box operations on a slide."""
        if operation == "add":
            slide, error = get_required_slide(presentation_id, slide_index)
            if error:
                return error

            try:
                ppt_utils.add_textbox(slide, left, top, width, height, text)
                return {
                    "message": f"Added text box to slide {slide_index}",
                    "shape_index": len(slide.shapes) - 1,
                    "text": text
                }
            except Exception as e:
                return {
                    "error": f"Failed to add text box: {str(e)}"
                }

        shape, error = get_required_shape(
            presentation_id, slide_index, shape_index)
        if error:
            return error

        if not hasattr(shape, 'text_frame') or not shape.text_frame:
            return {"error": "Shape does not contain text"}

        if operation == "update_text":
            try:
                shape.text = text
                return {
                    "message": f"Updated text shape {shape_index} on slide {slide_index}",
                    "shape_index": shape_index,
                    "text": text
                }
            except Exception as e:
                return {
                    "error": f"Failed to update text box text: {str(e)}"
                }

        if operation == "format":
            validations = {
                "font_size": (font_size, [(is_positive, "must be a positive integer")]),
                "color": (color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")]),
                "bg_color": (bg_color, [(is_valid_rgb, "must be a valid RGB list [R, G, B] with values 0-255")])
            }
            valid, error_message = validate_parameters(validations)
            if not valid:
                return {"error": error_message}

            try:
                ppt_utils.format_text_advanced(
                    shape,
                    font_size=font_size,
                    font_name=font_name,
                    bold=bold,
                    italic=italic,
                    underline=underline,
                    color=tuple(color),
                    bg_color=tuple(bg_color),
                    alignment=alignment,
                    vertical_alignment=vertical_alignment
                )
                return {
                    "message": f"Formatted text shape {shape_index} on slide {slide_index}"
                }
            except Exception as e:
                return {
                    "error": f"Failed to format text box: {str(e)}"
                }

        if operation == "set_runs":
            try:
                text_frame = shape.text_frame
                text_frame.clear()
                formatted_runs = []

                from pptx.util import Pt
                from pptx.dml.color import RGBColor

                for run_data in text_runs:
                    if 'text' not in run_data:
                        continue

                    paragraph = text_frame.paragraphs[0] if len(
                        formatted_runs) == 0 else text_frame.add_paragraph()
                    run = paragraph.add_run()
                    run.text = run_data['text']

                    if 'bold' in run_data:
                        run.font.bold = run_data['bold']
                    if 'italic' in run_data:
                        run.font.italic = run_data['italic']
                    if 'underline' in run_data:
                        run.font.underline = run_data['underline']
                    if 'font_size' in run_data:
                        run.font.size = Pt(run_data['font_size'])
                    if 'font_name' in run_data:
                        run.font.name = run_data['font_name']
                    if 'color' in run_data and is_valid_rgb(run_data['color']):
                        run.font.color.rgb = RGBColor(*run_data['color'])
                    if 'hyperlink' in run_data:
                        run.hyperlink.address = run_data['hyperlink']

                    formatted_runs.append({
                        "text": run_data['text'],
                        "formatting_applied": {k: v for k, v in run_data.items() if k != 'text'}
                    })

                return {
                    "message": f"Applied formatting to {len(formatted_runs)} text runs on shape {shape_index}",
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "formatted_runs": formatted_runs
                }
            except Exception as e:
                return {
                    "error": f"Failed to set text box runs: {str(e)}"
                }

        if operation == "validate_layout":
            slide, slide_error = get_required_slide(
                presentation_id, slide_index)
            if slide_error:
                return slide_error

            try:
                validation_result = ppt_utils.validate_text_fit(
                    shape,
                    text_content=None,
                    font_size=12
                )

                if validation_result.get("needs_optimization"):
                    fix_result = ppt_utils.validate_and_fix_slide(
                        slide,
                        auto_fix=True,
                        min_font_size=min_font_size,
                        max_font_size=max_font_size
                    )
                    validation_result.update(fix_result)

                return validation_result
            except Exception as e:
                return {
                    "error": f"Failed to validate text box layout: {str(e)}"
                }

        return {"error": "Invalid text operation. Use: add, update_text, format, set_runs, validate_layout"}

    @app.tool(
        annotations=ToolAnnotations(
            title="Manage Image",
        ),
    )
    def manage_image(
        presentation_id: str,
        operation: str,
        slide_index: int,
        image_source: str,
        left: float,
        top: float,
        width: float,
        height: float,
        enhancement_style: str,
        brightness: float,
        contrast: float,
        saturation: float,
        sharpness: float,
        blur_radius: float,
        filter_type: str,
        output_path: str,
    ) -> Dict:
        """Manage image insertion and enhancement operations."""
        if operation == "enhance":
            if not os.path.exists(image_source):
                return {
                    "error": f"Image file not found: {image_source}"
                }

            try:
                if enhancement_style == "presentation":
                    enhanced_path = ppt_utils.apply_professional_image_enhancement(
                        image_source, style="presentation", output_path=output_path
                    )
                else:
                    enhanced_path = ppt_utils.enhance_image_with_pillow(
                        image_source,
                        brightness=brightness,
                        contrast=contrast,
                        saturation=saturation,
                        sharpness=sharpness,
                        blur_radius=blur_radius,
                        filter_type=filter_type,
                        output_path=output_path
                    )

                return {
                    "message": f"Enhanced image: {image_source}",
                    "enhanced_path": enhanced_path
                }
            except Exception as e:
                return {
                    "error": f"Failed to enhance image: {str(e)}"
                }

        slide, error = get_required_slide(presentation_id, slide_index)
        if error:
            return error

        temp_path = None
        image_path = image_source

        try:
            if operation == "add_file":
                if not os.path.exists(image_source):
                    return {
                        "error": f"Image file not found: {image_source}"
                    }
            elif operation == "add_url":
                import urllib.request

                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                    temp_path = temp_file.name

                urllib.request.urlretrieve(image_source, temp_path)
                image_path = temp_path
            elif operation == "add_base64":
                image_data = base64.b64decode(image_source)
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                    temp_file.write(image_data)
                    temp_path = temp_file.name
                image_path = temp_path
            else:
                return {"error": "Invalid image operation. Use: add_file, add_url, add_base64, enhance"}

            ppt_utils.add_image(slide, image_path, left, top, width, height)
            return {
                "message": f"Added image to slide {slide_index}",
                "shape_index": len(slide.shapes) - 1,
                "image_source": image_source,
                "operation": operation
            }
        except Exception as e:
            if operation == "add_url":
                return {
                    "error": f"Failed to add image from URL: {str(e)}"
                }
            if operation == "add_base64":
                return {
                    "error": f"Failed to process base64 image: {str(e)}"
                }
            return {
                "error": f"Failed to add image from file: {str(e)}"
            }
        finally:
            if temp_path and os.path.exists(temp_path):
                os.unlink(temp_path)
