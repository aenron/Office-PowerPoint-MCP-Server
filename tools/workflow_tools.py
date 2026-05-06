"""
High-level workflow tools for template-aware PowerPoint generation.
These tools provide a simplified MCP surface while reusing the existing
lower-level PowerPoint utilities internally.
"""
from typing import Any, Dict, List, Optional, Tuple
import json
import os
import re

from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches
from starlette.requests import Request
from starlette.responses import FileResponse, JSONResponse

import utils as ppt_utils


def register_workflow_tools(
    app: FastMCP,
    presentations: Dict,
    projects: Dict,
    get_current_presentation_id,
    set_current_presentation_id,
    get_template_search_directories,
):
    """Register simplified, template-aware workflow tools."""
    project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    download_dir = os.path.join(project_root, "ppt")
    download_url = os.environ.get("DOWNLOAD_URL")
    emu_per_inch = 914400
    profile_dir = os.path.join(project_root, "templates", "profiles")
    os.makedirs(download_dir, exist_ok=True)
    os.makedirs(profile_dir, exist_ok=True)

    def normalize_text(value: str) -> str:
        return re.sub(r"[\s_\-]+", "", (value or "").strip().lower())

    def slugify(value: str) -> str:
        cleaned = re.sub(r"[^A-Za-z0-9\u4e00-\u9fff]+",
                         "-", (value or "").strip())
        cleaned = cleaned.strip("-")
        return cleaned or "presentation"

    def is_env_enabled(name: str) -> bool:
        return (os.environ.get(name) or "").strip().lower() in {"1", "true", "yes", "on"}

    def get_builtin_themes() -> Dict[str, Dict[str, Any]]:
        return {
            "business_blue": {
                "name": "Business Blue",
                "description": "Clean corporate style for proposals, reports, and data platform decks.",
                "font_name": "Microsoft YaHei",
                "colors": {
                    "background": (248, 250, 252),
                    "surface": (255, 255, 255),
                    "primary": (32, 45, 70),
                    "secondary": (69, 84, 112),
                    "accent": (20, 126, 214),
                    "muted": (107, 114, 128),
                    "line": (218, 224, 232),
                    "light": (235, 241, 248),
                    "danger": (185, 45, 55),
                    "success": (42, 132, 84),
                },
            },
            "modern_green": {
                "name": "Modern Green",
                "description": "Restrained green theme for transformation, sustainability, and operations.",
                "font_name": "Microsoft YaHei",
                "colors": {
                    "background": (247, 250, 248),
                    "surface": (255, 255, 255),
                    "primary": (26, 69, 57),
                    "secondary": (60, 91, 82),
                    "accent": (0, 145, 110),
                    "muted": (100, 116, 109),
                    "line": (214, 225, 220),
                    "light": (230, 242, 237),
                    "danger": (181, 57, 57),
                    "success": (0, 128, 90),
                },
            },
            "executive_gray": {
                "name": "Executive Gray",
                "description": "Neutral executive style for concise management communication.",
                "font_name": "Microsoft YaHei",
                "colors": {
                    "background": (249, 249, 247),
                    "surface": (255, 255, 255),
                    "primary": (38, 43, 50),
                    "secondary": (84, 91, 100),
                    "accent": (198, 111, 52),
                    "muted": (112, 117, 124),
                    "line": (222, 222, 218),
                    "light": (238, 238, 234),
                    "danger": (174, 65, 65),
                    "success": (68, 132, 94),
                },
            },
            "academic_burgundy": {
                "name": "Academic Burgundy",
                "description": "Restrained academic style for social-science talks, thesis defenses, and research seminars.",
                "font_name": "Microsoft YaHei",
                "colors": {
                    "background": (250, 250, 248),
                    "surface": (255, 255, 255),
                    "primary": (36, 42, 53),
                    "secondary": (82, 88, 99),
                    "accent": (136, 46, 67),
                    "muted": (116, 119, 126),
                    "line": (224, 224, 220),
                    "light": (242, 238, 242),
                    "danger": (170, 70, 78),
                    "success": (54, 125, 118),
                },
            },
            "academic_default": {
                "name": "学术默认主题",
                "description": "Academic default theme derived from templates/智算专家会.pptx with deep technology blue accents.",
                "font_name": "Microsoft YaHei",
                "colors": {
                    "background": (245, 247, 250),
                    "surface": (255, 255, 255),
                    "primary": (10, 42, 84),
                    "secondary": (80, 80, 80),
                    "accent": (74, 144, 226),
                    "muted": (110, 116, 122),
                    "line": (222, 224, 227),
                    "light": (232, 241, 251),
                    "danger": (237, 125, 49),
                    "success": (80, 150, 120),
                },
            },
        }

    def get_builtin_layouts() -> List[Dict[str, Any]]:
        return [
            {
                "layout_id": "cover",
                "name": "Cover",
                "description": "Opening slide with title, subtitle, and deck tag.",
                "use_when": "Use for the first page of a deck.",
                "required_fields": ["title"],
                "optional_fields": ["subtitle", "content", "tag"],
                "supported_fields": ["type", "title", "subtitle", "content", "tag"],
                "capacity_limits": {
                    "subtitle_chars": 80,
                    "content_lines": 3,
                    "content_chars_per_line": 32,
                },
            },
            {
                "layout_id": "summary",
                "name": "Summary",
                "description": "Overview slide with a statement and up to four structured sections.",
                "use_when": "Use for agenda, chapter overview, executive summary, or final summary.",
                "required_fields": ["title"],
                "optional_fields": ["statement", "subtitle", "sections", "items", "content", "text"],
                "supported_fields": ["type", "title", "subtitle", "statement", "sections", "items", "content", "text", "source_note"],
                "capacity_limits": {
                    "sections": 4,
                    "items": 6,
                    "line_chars": 36,
                },
                "auto_infer_from": ["items", "sections"],
            },
            {
                "layout_id": "cards",
                "name": "Cards",
                "description": "Card grid for multiple concepts, modules, capabilities, or arguments.",
                "use_when": "Use when the slide has 2-6 peer items.",
                "required_fields": ["title", "items"],
                "optional_fields": ["sections", "density"],
                "supported_fields": ["type", "title", "items", "sections", "density", "source_note"],
                "capacity_limits": {
                    "cards": 6,
                    "points_per_card": 4,
                    "line_chars": 32,
                },
                "auto_infer_from": ["items", "sections"],
            },
            {
                "layout_id": "comparison",
                "name": "Comparison",
                "description": "Two-column comparison for current state versus target state, or before versus after.",
                "use_when": "Use for contrasts, debates, alternatives, and problem-solution pages.",
                "required_fields": ["title"],
                "optional_fields": ["left", "right", "comparisons"],
                "supported_fields": ["type", "title", "left", "right", "comparisons", "source_note"],
                "capacity_limits": {
                    "columns": 2,
                    "items_per_column": 5,
                    "line_chars": 34,
                },
                "auto_infer_from": ["left", "right", "comparisons"],
            },
            {
                "layout_id": "process",
                "name": "Process",
                "description": "Horizontal process with up to four steps.",
                "use_when": "Use for method flow, implementation path, or logical sequence.",
                "required_fields": ["title", "steps"],
                "optional_fields": ["items", "sections"],
                "supported_fields": ["type", "title", "steps", "items", "sections", "source_note"],
                "capacity_limits": {
                    "steps": 4,
                    "points_per_step": 3,
                    "line_chars": 28,
                },
                "auto_infer_from": ["steps"],
            },
            {
                "layout_id": "timeline",
                "name": "Timeline",
                "description": "Timeline view for phased plans, research stages, or historical evolution.",
                "use_when": "Use when steps have phases or more than four sequential items.",
                "required_fields": ["title", "steps"],
                "optional_fields": ["items"],
                "supported_fields": ["type", "title", "steps", "items", "source_note"],
                "capacity_limits": {
                    "steps": 8,
                    "line_chars": 28,
                },
                "auto_infer_from": ["steps"],
            },
            {
                "layout_id": "metrics",
                "name": "Metrics",
                "description": "Metric cards for quantified indicators, sample statistics, or key numbers.",
                "use_when": "Use only when the input includes concrete numbers or indicators.",
                "required_fields": ["title", "metrics"],
                "optional_fields": ["items"],
                "supported_fields": ["type", "title", "metrics", "items", "source_note"],
                "capacity_limits": {
                    "metrics": 6,
                    "label_chars": 16,
                    "detail_chars": 32,
                },
                "auto_infer_from": ["metrics"],
            },
            {
                "layout_id": "architecture",
                "name": "Architecture",
                "description": "Layered structure for systems, governance structures, or analytical levels.",
                "use_when": "Use for multi-layer frameworks and component relationships.",
                "required_fields": ["title", "layers"],
                "optional_fields": ["items", "relations"],
                "supported_fields": ["type", "title", "layers", "items", "relations", "source_note"],
                "capacity_limits": {
                    "layers": 5,
                    "items_per_layer": 4,
                    "line_chars": 28,
                },
                "auto_infer_from": ["layers", "relations"],
            },
            {
                "layout_id": "table",
                "name": "Table",
                "description": "Structured table with headers and rows.",
                "use_when": "Use for variables, samples, coding schemes, literature details, and comparisons.",
                "required_fields": ["title", "table"],
                "optional_fields": ["table.headers", "table.rows"],
                "supported_fields": ["type", "title", "table", "source_note"],
                "capacity_limits": {
                    "columns": 5,
                    "rows_per_slide": 6,
                    "cell_chars": 28,
                },
                "auto_infer_from": ["table"],
            },
            {
                "layout_id": "quote",
                "name": "Quote",
                "description": "Large conclusion or core argument with optional supporting points.",
                "use_when": "Use for thesis statement, key conclusion, or central claim.",
                "required_fields": ["title", "statement"],
                "optional_fields": ["content", "points", "items"],
                "supported_fields": ["type", "title", "statement", "content", "points", "items", "source_note"],
                "capacity_limits": {
                    "statement_chars": 90,
                    "points": 4,
                    "point_chars": 32,
                },
                "auto_infer_from": ["statement"],
            },
            {
                "layout_id": "closing",
                "name": "Closing",
                "description": "Final thank-you slide.",
                "use_when": "Use for the last page.",
                "required_fields": ["title"],
                "optional_fields": ["subtitle", "content"],
                "supported_fields": ["type", "title", "subtitle", "content", "source_note"],
                "capacity_limits": {
                    "subtitle_chars": 60,
                    "content_lines": 3,
                },
            },
            {
                "layout_id": "research_questions",
                "name": "Research Questions",
                "description": "Academic research question slide with context, questions, and optional research gap.",
                "use_when": "Use for opening the research logic in social-science talks.",
                "required_fields": ["title", "questions"],
                "optional_fields": ["background", "context", "statement", "gap", "research_gap"],
                "supported_fields": ["type", "title", "questions", "research_questions", "items", "background", "context", "statement", "gap", "research_gap", "source_note"],
                "capacity_limits": {
                    "questions": 4,
                    "background_chars": 70,
                    "context_chars": 90,
                    "gap_chars": 90,
                },
                "auto_infer_from": ["questions", "research_questions"],
            },
            {
                "layout_id": "literature_matrix",
                "name": "Literature Matrix",
                "description": "Literature review matrix covering research stream, main argument, gap, and this study's entry point.",
                "use_when": "Use for literature review and positioning pages.",
                "required_fields": ["title", "literature"],
                "optional_fields": ["studies", "items", "table.headers", "table.rows"],
                "supported_fields": ["type", "title", "literature", "studies", "items", "table", "source_note"],
                "capacity_limits": {
                    "rows_per_slide": 6,
                    "columns": 4,
                    "cell_chars": 30,
                },
                "auto_infer_from": ["literature", "studies"],
            },
            {
                "layout_id": "theoretical_framework",
                "name": "Theoretical Framework",
                "description": "Concept or variable relationship page with propositions and mechanism explanation.",
                "use_when": "Use for theory, hypotheses, analytical framework, and mechanism pages.",
                "required_fields": ["title", "concepts"],
                "optional_fields": ["framework", "variables", "relations", "propositions", "hypotheses", "mechanism", "explanation"],
                "supported_fields": ["type", "title", "concepts", "framework", "variables", "items", "relations", "propositions", "hypotheses", "mechanism", "explanation", "content", "source_note"],
                "capacity_limits": {
                    "concepts": 5,
                    "concept_title_chars": 14,
                    "points_per_concept": 2,
                    "point_chars": 24,
                    "framework_chars": 120,
                },
                "auto_infer_from": ["framework", "concepts", "variables"],
            },
            {
                "layout_id": "method_design",
                "name": "Method Design",
                "description": "Research design page covering data, sample, variables, and analysis methods.",
                "use_when": "Use for data and method sections in academic reports.",
                "required_fields": ["title"],
                "optional_fields": ["methods", "research_design", "data_sources", "sample", "variables", "analysis", "method"],
                "supported_fields": ["type", "title", "methods", "research_design", "steps", "data_sources", "data", "sample", "scope", "variables", "analysis", "method", "source_note"],
                "capacity_limits": {
                    "sections": 4,
                    "points_per_section": 3,
                    "line_chars": 34,
                },
                "auto_infer_from": ["methods", "data_sources", "research_design"],
            },
            {
                "layout_id": "findings",
                "name": "Findings",
                "description": "Main findings page with concise claims and supporting evidence.",
                "use_when": "Use for empirical findings, analysis results, or case observations.",
                "required_fields": ["title", "findings"],
                "optional_fields": ["headline", "statement", "items", "sections"],
                "supported_fields": ["type", "title", "findings", "headline", "statement", "items", "sections", "source_note"],
                "capacity_limits": {
                    "findings": 4,
                    "points_per_finding": 3,
                    "point_chars": 28,
                },
                "auto_infer_from": ["findings"],
            },
            {
                "layout_id": "contribution_limitations",
                "name": "Contribution And Limitations",
                "description": "Academic contribution, limitations, and future research directions.",
                "use_when": "Use near the end of a research presentation.",
                "required_fields": ["title"],
                "optional_fields": ["contributions", "limitations", "implications", "future", "outlook"],
                "supported_fields": ["type", "title", "contributions", "contribution", "limitations", "limitation", "implications", "future", "outlook", "left", "right", "items", "source_note"],
                "capacity_limits": {
                    "contributions": 3,
                    "limitations": 3,
                    "future": 3,
                    "line_chars": 30,
                },
                "auto_infer_from": ["contributions", "limitations"],
            },
            {
                "layout_id": "expert_section",
                "name": "智算章节页",
                "description": "Section divider inspired by templates/智算专家会.pptx, with centered title and subtitle.",
                "use_when": "Use for chapter openings and topic transitions.",
                "required_fields": ["title"],
                "optional_fields": ["subtitle", "content", "tag"],
                "supported_fields": ["type", "title", "subtitle", "content", "tag", "source_note"],
                "capacity_limits": {
                    "title_chars": 26,
                    "subtitle_chars": 42,
                    "content_lines": 2,
                },
            },
            {
                "layout_id": "expert_title_content",
                "name": "智算标题内容页",
                "description": "Top-centered title layout with open content blocks, matching the expert forum deck rhythm.",
                "use_when": "Use for overview, explanation, and normal content pages.",
                "required_fields": ["title"],
                "optional_fields": ["statement", "items", "sections", "content"],
                "supported_fields": ["type", "title", "statement", "items", "sections", "content", "source_note"],
                "capacity_limits": {
                    "blocks": 4,
                    "points_per_block": 3,
                    "line_chars": 32,
                },
            },
            {
                "layout_id": "expert_split",
                "name": "智算左右图文页",
                "description": "Large left visual area plus right explanatory stack, based on the case-study pages in the template.",
                "use_when": "Use for case pages, scenario explanation, and image-plus-text evidence.",
                "required_fields": ["title"],
                "optional_fields": ["image_path", "image_caption", "items", "sections", "content"],
                "supported_fields": ["type", "title", "image_path", "image_caption", "items", "sections", "content", "source_note"],
                "capacity_limits": {
                    "right_blocks": 4,
                    "points_per_block": 2,
                    "line_chars": 30,
                },
            },
            {
                "layout_id": "expert_path",
                "name": "智算流程路径页",
                "description": "Six-step implementation path with a central numbered rail and side explanations.",
                "use_when": "Use for scenario implementation paths, process methods, or staged logic.",
                "required_fields": ["title", "steps"],
                "optional_fields": ["items", "sections"],
                "supported_fields": ["type", "title", "steps", "items", "sections", "source_note"],
                "capacity_limits": {
                    "steps": 6,
                    "step_title_chars": 18,
                    "step_detail_chars": 36,
                },
            },
            {
                "layout_id": "expert_scope",
                "name": "智算范围清单页",
                "description": "Vertical scope list inspired by the security-scope slide in the template.",
                "use_when": "Use for scope, safeguards, checklist, policy coverage, or risk categories.",
                "required_fields": ["title", "items"],
                "optional_fields": ["sections"],
                "supported_fields": ["type", "title", "items", "sections", "source_note"],
                "capacity_limits": {
                    "items": 4,
                    "points_per_item": 2,
                    "line_chars": 32,
                },
            },
        ]

    def get_theme(theme_id: str) -> Tuple[str, Dict[str, Any]]:
        themes = get_builtin_themes()
        effective_theme_id = theme_id if theme_id in themes else "business_blue"
        return effective_theme_id, themes[effective_theme_id]

    def rgb(value: Tuple[int, int, int]) -> RGBColor:
        return RGBColor(value[0], value[1], value[2])

    def theme_color(theme: Dict[str, Any], role: str) -> Tuple[int, int, int]:
        return theme["colors"].get(role, theme["colors"]["primary"])

    def fit_text_size(text: str, base_size: int, min_size: int, soft_limit: int) -> int:
        length = len(text or "")
        if length <= soft_limit:
            return base_size
        reduced = base_size - int((length - soft_limit) / 12) - 1
        return max(min_size, reduced)

    def density_factor(density: str) -> float:
        return {
            "spacious": 0.82,
            "standard": 1.0,
            "compact": 1.28,
        }.get((density or "standard").strip().lower(), 1.0)

    def text_capacity(width: float, height: float, font_size: int, density: str = "standard") -> int:
        base = max(12, int(width * height * 18 * density_factor(density)))
        if font_size <= 9:
            return int(base * 1.35)
        if font_size >= 18:
            return int(base * 0.62)
        return base

    def adapt_text(
        text: Any,
        width: float,
        height: float,
        base_size: int,
        min_size: int,
        density: str,
        overflow: str,
        warnings: Optional[List[str]],
        context: str,
    ) -> Tuple[str, int]:
        value = str(text or "").strip()
        capacity = text_capacity(width, height, base_size, density)
        size = fit_text_size(value, base_size, min_size, capacity)
        effective_overflow = (
            overflow or "shrink_then_truncate").strip().lower()

        if len(value) <= capacity:
            return value, size

        if "truncate" in effective_overflow:
            max_chars = max(8, int(capacity * (base_size / max(size, 1))))
            if len(value) > max_chars:
                value = value[: max_chars - 1].rstrip() + "..."
                if warnings is not None:
                    warnings.append(
                        f"{context} truncated from {len(str(text or ''))} to {len(value)} characters.")
        elif "warn" in effective_overflow and warnings is not None:
            warnings.append(
                f"{context} may overflow: {len(value)} characters for estimated capacity {capacity}.")

        return value, size

    def safe_lines(values: Any, limit: int = 4) -> List[str]:
        if values is None:
            return []
        if isinstance(values, str):
            return [line.strip() for line in re.split(r"[\n；;]+", values) if line.strip()][:limit]
        if isinstance(values, list):
            lines = []
            for item in values:
                if isinstance(item, dict):
                    text = item.get("text") or item.get("title") or item.get(
                        "label") or item.get("name") or ""
                    detail = item.get("detail") or item.get(
                        "description") or item.get("content") or ""
                    lines.append(f"{text}: {detail}" if detail else str(text))
                else:
                    lines.append(str(item))
            return [line.strip() for line in lines if line and line.strip()][:limit]
        return [str(values)][:limit]

    def split_label_detail(value: Any) -> Tuple[str, str]:
        text = str(value or "").strip()
        if not text:
            return "", ""
        match = re.match(r"^([^:：]{1,12})[:：]\s*(.+)$", text)
        if not match:
            return text, ""
        return match.group(1).strip(), match.group(2).strip()

    def source_refs_to_note(value: Any) -> str:
        if not value:
            return ""
        if isinstance(value, str):
            return value.strip()
        if isinstance(value, dict):
            parts = [
                str(value.get(key)).strip()
                for key in ["file", "section", "chapter", "page", "table", "paragraph", "quote"]
                if value.get(key)
            ]
            return " / ".join(parts)
        if isinstance(value, list):
            notes = [source_refs_to_note(item) for item in value[:3]]
            return "；".join([item for item in notes if item])
        return str(value).strip()

    def get_source_note(item: Dict[str, Any]) -> str:
        note = (
            item.get("source_note")
            or item.get("source")
            or item.get("source_text")
            or item.get("citation")
            or item.get("reference")
            or source_refs_to_note(item.get("source_refs"))
        )
        return str(note or "").strip()

    def merge_point_lines(item: Dict[str, Any], fallback: Any = "") -> List[str]:
        lines: List[str] = []
        for key in ["points", "bullets", "content", "details", "evidence", "explanation", "analysis", "result", "conclusion", "mechanism", "boundary"]:
            for line in safe_lines(item.get(key), 6):
                if line and line not in lines:
                    lines.append(line)
        for line in safe_lines(fallback, 6):
            if line and line not in lines:
                lines.append(line)
        return lines

    def trim_items(items: List[Any], limit: int, warnings: Optional[List[str]], context: str) -> List[Any]:
        if len(items) > limit and warnings is not None:
            warnings.append(
                f"{context} limited to {limit} items; {len(items) - limit} extra items were not rendered.")
        return items[:limit]

    def add_text(
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        theme: Dict[str, Any],
        font_size: int,
        color_role: str = "primary",
        bold: bool = False,
        alignment: str = "left",
        vertical_alignment: str = "top",
        density: str = "standard",
        overflow: str = "shrink_then_truncate",
        min_font_size: int = 8,
        warnings: Optional[List[str]] = None,
        context: str = "text",
    ):
        text, font_size = adapt_text(
            text,
            width,
            height,
            font_size,
            min_font_size,
            density,
            overflow,
            warnings,
            context,
        )
        shape = ppt_utils.add_textbox(
            slide,
            left,
            top,
            width,
            height,
            text or "",
            font_size=font_size,
            font_name=theme.get("font_name"),
            bold=bold,
            color=theme_color(theme, color_role),
            alignment=alignment,
            vertical_alignment=vertical_alignment,
        )
        shape.text_frame.margin_left = Inches(0.04)
        shape.text_frame.margin_right = Inches(0.04)
        shape.text_frame.margin_top = Inches(0.02)
        shape.text_frame.margin_bottom = Inches(0.02)
        return shape

    def add_rect(
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        fill: Tuple[int, int, int],
        line: Optional[Tuple[int, int, int]] = None,
        radius: bool = False,
    ):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, Inches(
            left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        if line:
            shape.line.color.rgb = rgb(line)
            shape.line.width = Inches(0.01)
        else:
            shape.line.fill.background()
        return shape

    def add_theme_background(slide, theme: Dict[str, Any]) -> None:
        add_rect(slide, 0, 0, 13.333, 7.5, theme_color(theme, "background"))
        add_rect(slide, 0, 0, 0.18, 7.5, theme_color(theme, "accent"))
        add_rect(slide, 0.18, 0, 13.153, 0.08, theme_color(theme, "primary"))

    def add_slide_header(slide, title: str, theme: Dict[str, Any], kicker: str = "") -> None:
        add_theme_background(slide, theme)
        if kicker:
            add_text(slide, 0.7, 0.35, 3.4, 0.28,
                     kicker, theme, 8, "accent", True)
        add_text(
            slide,
            0.7,
            0.58,
            8.8,
            0.55,
            title,
            theme,
            fit_text_size(title, 24, 17, 24),
            "primary",
            True,
        )
        add_rect(slide, 0.72, 1.28, 1.1, 0.04, theme_color(theme, "accent"))

    def add_source_note(slide, slide_spec: Dict[str, Any], theme: Dict[str, Any], warnings: Optional[List[str]] = None) -> None:
        note = get_source_note(slide_spec)
        if not note:
            return
        label = note if note.startswith("来源") else f"来源：{note}"
        add_text(
            slide,
            0.72,
            6.72,
            9.8,
            0.22,
            label,
            theme,
            6,
            "muted",
            density="compact",
            overflow="shrink_then_truncate",
            min_font_size=6,
            warnings=warnings,
            context="source note",
        )

    def apply_page_size(presentation, page_size: str) -> str:
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        return "wide_16_9"

    def is_cover_spec(slide_spec: Dict[str, Any]) -> bool:
        return infer_slide_type(slide_spec) == "cover"

    def is_closing_spec(slide_spec: Dict[str, Any]) -> bool:
        return infer_slide_type(slide_spec) == "closing"

    def prepare_slide_specs(
        slides: List[Dict[str, Any]],
        title: str,
        subtitle: str,
        auto_cover: bool,
        auto_closing: bool,
        style: str,
    ) -> List[Dict[str, Any]]:
        prepared = list(slides or [])
        if not prepared:
            prepared = [
                {
                    "type": "summary",
                    "title": title,
                    "statement": subtitle,
                    "sections": [{"title": "核心观点", "points": safe_lines(subtitle, 3)}],
                }
            ]
        if auto_cover and (not prepared or not is_cover_spec(prepared[0] if isinstance(prepared[0], dict) else {})):
            prepared.insert(0, {
                "type": "cover",
                "title": title,
                "subtitle": subtitle,
                "tag": style.upper() if style else "PPT",
            })
        if auto_closing and (not prepared or not is_closing_spec(prepared[-1] if isinstance(prepared[-1], dict) else {})):
            prepared.append({
                "type": "closing",
                "title": "谢谢",
                "subtitle": title,
            })
        return prepared

    def chunk_items(items: List[Any], size: int) -> List[List[Any]]:
        if size <= 0:
            return [items]
        return [items[index:index + size] for index in range(0, len(items), size)] or [[]]

    def part_title(title: str, index: int, total: int) -> str:
        if total <= 1:
            return title
        return f"{title}（{index}/{total}）"

    def clone_slide_part(slide_spec: Dict[str, Any], index: int, total: int) -> Dict[str, Any]:
        part = dict(slide_spec)
        part["title"] = part_title(str(slide_spec.get("title") or "内容"), index, total)
        return part

    def split_table_slide(slide_spec: Dict[str, Any]) -> List[Dict[str, Any]]:
        table_spec = slide_spec.get("table") or {}
        rows = table_spec.get("rows") or []
        if not isinstance(rows, list) or len(rows) <= 6:
            return [slide_spec]

        row_chunks = chunk_items(rows, 6)
        slides_out: List[Dict[str, Any]] = []
        for index, row_chunk in enumerate(row_chunks, start=1):
            part = clone_slide_part(slide_spec, index, len(row_chunks))
            part_table = dict(table_spec)
            part_table["rows"] = row_chunk
            part["table"] = part_table
            slides_out.append(part)
        return slides_out

    def split_list_slide(slide_spec: Dict[str, Any], field_name: str, limit: int) -> List[Dict[str, Any]]:
        values = slide_spec.get(field_name) or []
        if not isinstance(values, list) or len(values) <= limit:
            return [slide_spec]

        value_chunks = chunk_items(values, limit)
        slides_out: List[Dict[str, Any]] = []
        for index, value_chunk in enumerate(value_chunks, start=1):
            part = clone_slide_part(slide_spec, index, len(value_chunks))
            part[field_name] = value_chunk
            slides_out.append(part)
        return slides_out

    def expand_finding_items(findings: List[Any]) -> List[Any]:
        expanded: List[Any] = []
        for finding in findings:
            if not isinstance(finding, dict):
                expanded.append(finding)
                continue
            points = merge_point_lines(finding)
            if len(points) <= 3:
                expanded.append(finding)
                continue
            point_chunks = chunk_items(points, 3)
            for index, point_chunk in enumerate(point_chunks, start=1):
                part = dict(finding)
                part["points"] = point_chunk
                if index > 1:
                    part["title"] = f"{finding.get('title') or finding.get('headline') or '发现'}（续）"
                expanded.append(part)
        return expanded

    def split_findings_slide(slide_spec: Dict[str, Any]) -> List[Dict[str, Any]]:
        findings = slide_spec.get("findings") or slide_spec.get("items") or slide_spec.get("sections") or []
        if not isinstance(findings, list):
            return [slide_spec]

        expanded = expand_finding_items(findings)
        if len(expanded) <= 4:
            part = dict(slide_spec)
            if slide_spec.get("findings") is not None:
                part["findings"] = expanded
            elif slide_spec.get("sections") is not None:
                part["sections"] = expanded
            else:
                part["items"] = expanded
            return [part]

        finding_chunks = chunk_items(expanded, 4)
        slides_out: List[Dict[str, Any]] = []
        for index, finding_chunk in enumerate(finding_chunks, start=1):
            part = clone_slide_part(slide_spec, index, len(finding_chunks))
            if slide_spec.get("findings") is not None:
                part["findings"] = finding_chunk
            elif slide_spec.get("sections") is not None:
                part["sections"] = finding_chunk
            else:
                part["items"] = finding_chunk
            slides_out.append(part)
        return slides_out

    def split_contribution_limitations_slide(slide_spec: Dict[str, Any]) -> List[Dict[str, Any]]:
        field_names = ["contributions", "limitations", "implications", "future", "outlook"]
        chunks_by_field: Dict[str, List[List[Any]]] = {}
        max_parts = 1
        for field_name in field_names:
            values = safe_lines(slide_spec.get(field_name), 12)
            if not values:
                continue
            chunks = chunk_items(values, 3)
            chunks_by_field[field_name] = chunks
            max_parts = max(max_parts, len(chunks))
        if max_parts <= 1:
            return [slide_spec]

        slides_out: List[Dict[str, Any]] = []
        for index in range(max_parts):
            part = clone_slide_part(slide_spec, index + 1, max_parts)
            for field_name, chunks in chunks_by_field.items():
                part[field_name] = chunks[index] if index < len(chunks) else []
            slides_out.append(part)
        return slides_out

    def expand_capacity_slide_specs(slide_specs: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        expanded: List[Dict[str, Any]] = []
        for slide_spec in slide_specs:
            if not isinstance(slide_spec, dict):
                expanded.append(slide_spec)
                continue

            slide_type = infer_slide_type(slide_spec)
            split_specs = [slide_spec]
            if slide_type == "table":
                split_specs = split_table_slide(slide_spec)
            elif slide_type == "literature_matrix":
                field_name = "literature" if slide_spec.get("literature") is not None else ("studies" if slide_spec.get("studies") is not None else "items")
                split_specs = split_list_slide(slide_spec, field_name, 6)
            elif slide_type in {"cards", "section"}:
                field_name = "sections" if slide_spec.get("sections") is not None else "items"
                split_specs = split_list_slide(slide_spec, field_name, 6)
            elif slide_type == "findings":
                split_specs = split_findings_slide(slide_spec)
            elif slide_type == "contribution_limitations":
                split_specs = split_contribution_limitations_slide(slide_spec)

            expanded.extend(split_specs)
        return expanded

    def add_deck_footer(
        presentation,
        theme: Dict[str, Any],
        footer_text: str,
        show_footer: bool,
        show_page_number: bool,
        visual_level: str,
    ) -> None:
        if not show_footer and not show_page_number and visual_level == "clean":
            return
        total = len(presentation.slides)
        for index, slide in enumerate(presentation.slides, start=1):
            if visual_level in {"rich", "dense"}:
                add_rect(slide, 0.68, 7.08, 0.7, 0.04,
                         theme_color(theme, "accent"))
                if visual_level == "rich":
                    add_rect(slide, 11.78, 0.48, 0.28, 0.28,
                             theme_color(theme, "light"), radius=True)
                    add_rect(slide, 12.14, 0.48, 0.18, 0.18,
                             theme_color(theme, "accent"), radius=True)
            if show_footer and footer_text:
                add_text(slide, 0.72, 7.0, 6.5, 0.22,
                         footer_text, theme, 7, "muted")
            if show_page_number:
                add_text(slide, 11.75, 7.0, 0.85, 0.22,
                         f"{index:02d}/{total:02d}", theme, 7, "muted", alignment="right")

    def inspect_presentation_quality(presentation, warnings: List[str]) -> Dict[str, Any]:
        min_font_size: Optional[float] = None
        text_shape_count = 0
        empty_slide_count = 0
        for slide in presentation.slides:
            slide_has_text = False
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame") or not shape.text_frame:
                    continue
                if shape.text.strip():
                    text_shape_count += 1
                    slide_has_text = True
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            size = float(run.font.size.pt)
                            min_font_size = size if min_font_size is None else min(
                                min_font_size, size)
            if not slide_has_text:
                empty_slide_count += 1

        return {
            "warning_count": len(warnings),
            "truncated_count": sum(1 for item in warnings if "truncated" in item),
            "limited_item_count": sum(1 for item in warnings if "limited to" in item),
            "empty_slide_count": empty_slide_count,
            "text_shape_count": text_shape_count,
            "min_font_size": min_font_size,
        }

    def add_card(
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        title: str,
        body: List[str],
        theme: Dict[str, Any],
        accent_role: str = "accent",
        density: str = "standard",
        overflow: str = "shrink_then_truncate",
        warnings: Optional[List[str]] = None,
    ) -> None:
        add_rect(slide, left, top, width, height, theme_color(
            theme, "surface"), theme_color(theme, "line"), True)
        add_rect(slide, left, top, 0.08, height,
                 theme_color(theme, accent_role))
        add_text(
            slide,
            left + 0.22,
            top + 0.18,
            width - 0.38,
            0.38,
            title,
            theme,
            fit_text_size(title, 13, 10, 16),
            "primary",
            True,
            density=density,
            overflow=overflow,
            warnings=warnings,
            context=f"card title '{title}'",
        )
        body_limit = 5 if density == "compact" else 4
        lines = trim_items([str(line) for line in body],
                           body_limit, warnings, f"card '{title}'")
        body_text = "\n".join([f"- {line}" for line in lines])
        add_text(
            slide,
            left + 0.22,
            top + 0.68,
            width - 0.42,
            height - 0.82,
            body_text,
            theme,
            10 if density != "spacious" else 11,
            "secondary",
            density=density,
            overflow=overflow,
            warnings=warnings,
            context=f"card body '{title}'",
        )

    def make_blank_slide(presentation):
        layout_index = 6 if len(presentation.slide_layouts) > 6 else 0
        return ppt_utils.add_slide(presentation, layout_index)[0]

    def normalized_sections(slide_spec: Dict[str, Any]) -> List[Dict[str, Any]]:
        sections = slide_spec.get("sections") or slide_spec.get("items") or []
        normalized: List[Dict[str, Any]] = []
        if isinstance(sections, list):
            for item in sections:
                if isinstance(item, dict):
                    title = item.get("title") or item.get("headline") or item.get(
                        "label") or item.get("name") or item.get("finding") or item.get("result") or item.get("conclusion") or ""
                    description = item.get("description") or item.get(
                        "detail") or item.get("summary") or item.get("evidence") or item.get("explanation") or ""
                    points = merge_point_lines(item, description)
                    if description and points and description not in points:
                        points = [str(description)] + points
                    normalized.append({
                        "title": str(title),
                        "points": points,
                        "source_note": get_source_note(item),
                        "raw": item,
                    })
                else:
                    title, detail = split_label_detail(item)
                    normalized.append({
                        "title": title,
                        "points": safe_lines(detail, 6),
                        "source_note": "",
                        "raw": item,
                    })
        if not normalized:
            normalized.append({
                "title": slide_spec.get("title") or "核心观点",
                "points": merge_point_lines(slide_spec, slide_spec.get("text")),
                "source_note": get_source_note(slide_spec),
                "raw": {},
            })
        return normalized

    def normalize_slide_spec(slide_spec: Dict[str, Any], deck_title: str = "", deck_subtitle: str = "") -> Dict[str, Any]:
        spec = dict(slide_spec or {})
        spec.setdefault("title", deck_title or "演示文稿")
        if deck_subtitle and not spec.get("subtitle"):
            spec["subtitle"] = deck_subtitle
        if spec.get("body") and not spec.get("content"):
            spec["content"] = spec.get("body")
        if spec.get("bullets") and not spec.get("items"):
            spec["items"] = spec.get("bullets")
        if not spec.get("source_note"):
            source_notes = [get_source_note(spec)]
            for key in ["items", "sections", "findings", "steps", "concepts", "methods", "literature", "studies"]:
                values = spec.get(key) or []
                if isinstance(values, dict):
                    values = [values]
                if isinstance(values, list):
                    source_notes.extend([get_source_note(item) for item in values if isinstance(item, dict)])
            spec["source_note"] = "；".join([note for note in source_notes if note][:3])
        return spec

    def infer_slide_type(slide_spec: Dict[str, Any]) -> str:
        explicit = (slide_spec.get("slide_type")
                    or slide_spec.get("type") or "").strip().lower()
        if explicit:
            aliases = {
                "two_column": "comparison",
                "roadmap": "process",
                "end": "closing",
                "bullet": "summary",
                "literature": "literature_matrix",
                "research_question": "research_questions",
                "framework": "theoretical_framework",
                "method": "method_design",
                "finding": "findings",
                "contribution": "contribution_limitations",
                "academic_default_section": "expert_section",
                "academic_default_content": "expert_title_content",
                "academic_default_split": "expert_split",
                "academic_default_path": "expert_path",
                "academic_default_scope": "expert_scope",
            }
            return aliases.get(explicit, explicit)
        if slide_spec.get("questions") or slide_spec.get("research_questions"):
            return "research_questions"
        if slide_spec.get("literature") or slide_spec.get("studies"):
            return "literature_matrix"
        if slide_spec.get("framework") or slide_spec.get("concepts") or slide_spec.get("variables"):
            return "theoretical_framework"
        if slide_spec.get("methods") or slide_spec.get("data_sources") or slide_spec.get("research_design"):
            return "method_design"
        if slide_spec.get("findings"):
            return "findings"
        if slide_spec.get("contributions") or slide_spec.get("limitations"):
            return "contribution_limitations"
        if slide_spec.get("layers") or slide_spec.get("relations"):
            return "architecture"
        if slide_spec.get("metrics"):
            return "metrics"
        if slide_spec.get("table"):
            return "table"
        if slide_spec.get("steps"):
            return "timeline" if len(slide_spec.get("steps") or []) > 4 else "process"
        if slide_spec.get("left") or slide_spec.get("right") or slide_spec.get("comparisons"):
            return "comparison"
        if slide_spec.get("statement") and not slide_spec.get("items") and not slide_spec.get("sections"):
            return "quote"
        if len(slide_spec.get("items") or slide_spec.get("sections") or []) > 0:
            return "cards"
        return "summary"

    def layout_metadata(layout_id: str) -> Dict[str, Any]:
        for layout in get_builtin_layouts():
            if layout.get("layout_id") == layout_id:
                return layout
        return {}

    def collect_ignored_fields(slide_spec: Dict[str, Any], rendered_type: str, slide_index: int) -> List[Dict[str, Any]]:
        layout = layout_metadata(rendered_type)
        supported = set(layout.get("supported_fields") or [])
        if not supported:
            return []
        ignored = []
        common_fields = {
            "slide_type",
            "kicker",
            "density",
            "overflow",
            "speaker_notes",
            "notes",
        }
        for key, value in (slide_spec or {}).items():
            if key in supported or key in common_fields:
                continue
            if value in (None, "", [], {}):
                continue
            ignored.append({
                "slide_index": slide_index,
                "slide_type": rendered_type,
                "field": key,
                "message": f"{rendered_type}.{key} is not rendered by the selected layout.",
            })
        return ignored

    def render_cover_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_theme_background(slide, theme)
        add_rect(slide, 0.75, 1.35, 1.5, 0.08, theme_color(theme, "accent"))
        title = slide_spec.get("title") or "演示文稿"
        subtitle = slide_spec.get(
            "subtitle") or slide_spec.get("content") or ""
        add_text(slide, 0.78, 1.65, 8.2, 1.2, title, theme, fit_text_size(title, 34, 24, 18), "primary", True,
                 density=density, overflow=overflow, min_font_size=20, warnings=warnings, context="cover title")
        add_text(slide, 0.82, 3.0, 7.8, 0.8, subtitle, theme, 15, "secondary", density=density,
                 overflow=overflow, min_font_size=10, warnings=warnings, context="cover subtitle")
        if slide_spec.get("content") and slide_spec.get("subtitle"):
            add_text(slide, 0.82, 4.12, 7.8, 0.9, slide_spec.get("content"), theme, 10, "muted",
                     density=density, overflow=overflow, min_font_size=8, warnings=warnings, context="cover content")
        add_rect(slide, 9.6, 0.0, 3.7, 7.5, theme_color(theme, "primary"))
        add_rect(slide, 10.1, 1.15, 2.35, 0.72,
                 theme_color(theme, "accent"), radius=True)
        add_text(slide, 10.36, 1.3, 1.85, 0.32, slide_spec.get(
            "tag") or "PPT", theme, 13, "surface", True, "center")
        for idx, y in enumerate([2.45, 3.15, 3.85]):
            add_rect(slide, 10.15 + idx * 0.25, y, 1.8,
                     0.08, theme_color(theme, "light"))

    def render_summary_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "方案概览",
                         theme, slide_spec.get("kicker") or "KEY POINTS")
        statement = slide_spec.get(
            "statement") or slide_spec.get("subtitle") or ""
        if statement:
            add_text(slide, 0.78, 1.45, 11.8, 0.42, statement, theme, 12, "secondary",
                     density=density, overflow=overflow, warnings=warnings, context="summary statement")
        sections = trim_items(normalized_sections(
            slide_spec), 4, warnings, "summary sections")
        positions = [(0.78, 1.7), (6.95, 1.7), (0.78, 4.35), (6.95, 4.35)]
        for index, section in enumerate(sections):
            left, top = positions[index]
            add_card(slide, left, top + (0.25 if statement else 0), 5.65, 2.0 if statement else 2.12,
                     section["title"], section["points"], theme, density=density, overflow=overflow, warnings=warnings)

    def render_cards_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "核心内容",
                         theme, slide_spec.get("kicker") or "OVERVIEW")
        sections = trim_items(normalized_sections(
            slide_spec), 6, warnings, "cards")
        if len(sections) <= 3:
            width, height = 3.75, 4.45
            for index, section in enumerate(sections):
                add_card(slide, 0.85 + index * 4.15, 1.8, width, height,
                         section["title"], section["points"], theme, density=density, overflow=overflow, warnings=warnings)
        else:
            width, height = 3.75, 2.05
            for index, section in enumerate(sections):
                row, col = divmod(index, 3)
                add_card(slide, 0.85 + col * 4.15, 1.65 + row * 2.48, width, height,
                         section["title"], section["points"], theme, density=density, overflow=overflow, warnings=warnings)

    def render_two_column_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "对比分析",
                         theme, slide_spec.get("kicker") or "COMPARISON")
        left = slide_spec.get("left") or {}
        right = slide_spec.get("right") or {}
        comparisons = slide_spec.get("comparisons") or []
        if comparisons and not (left or right):
            comparison_items = [
                item for item in comparisons if isinstance(item, dict)]
            column_items = [
                item for item in comparison_items
                if not any(item.get(key) for key in ["before", "after", "left", "right", "result"])
                and (item.get("title") or item.get("content") or item.get("points"))
            ]
            if len(column_items) >= 2 and len(column_items) == len(comparison_items):
                left = {
                    "title": column_items[0].get("title") or column_items[0].get("name") or "现状/痛点",
                    "points": safe_lines(column_items[0].get("points") or column_items[0].get("content") or column_items[0].get("description"), 8),
                }
                right = {
                    "title": column_items[1].get("title") or column_items[1].get("name") or "目标/方案",
                    "points": safe_lines(column_items[1].get("points") or column_items[1].get("content") or column_items[1].get("description"), 8),
                }
            else:
                left = {"title": "对比项", "points": [item.get("before") or item.get("left") or item.get(
                    "name") or item.get("title") or "" for item in comparison_items]}
                right = {"title": "优化后", "points": [item.get("after") or item.get("right") or item.get(
                    "result") or item.get("content") or "" for item in comparison_items]}
        add_card(slide, 0.8, 1.75, 5.65, 4.85, left.get("title") or "现状/痛点", safe_lines(left.get(
            "points") or left.get("content"), 8), theme, "danger", density, overflow, warnings)
        add_card(slide, 6.85, 1.75, 5.65, 4.85, right.get("title") or "目标/方案", safe_lines(right.get(
            "points") or right.get("content"), 8), theme, "success", density, overflow, warnings)

    def render_process_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "推进路径",
                         theme, slide_spec.get("kicker") or "ROADMAP")
        steps_source = slide_spec.get("steps") or slide_spec.get(
            "items") or slide_spec.get("sections")
        steps = trim_items(normalized_sections(
            {"items": steps_source or normalized_sections(slide_spec)}), 4, warnings, "process steps")
        step_width = 2.65
        y = 3.05
        for index, step in enumerate(steps):
            left = 0.82 + index * 3.05
            add_rect(slide, left, y, step_width, 0.78, theme_color(
                theme, "primary") if index % 2 == 0 else theme_color(theme, "secondary"))
            add_text(slide, left + 0.12, y + 0.18, 0.6, 0.3,
                     f"{index + 1:02d}", theme, 16, "surface", True, "center")
            add_text(slide, left + 0.82, y + 0.15, 1.68, 0.35, step["title"], theme, fit_text_size(
                step["title"], 12, 9, 9), "surface", True, density=density, overflow=overflow, warnings=warnings, context="process step title")
            add_text(slide, left, y + 1.05, step_width, 1.45, "\n".join(step["points"][:3]), theme, 10, "secondary",
                     alignment="center", density=density, overflow=overflow, warnings=warnings, context="process step body")
            if index < len(steps) - 1:
                add_text(slide, left + step_width + 0.16, y + 0.2, 0.28,
                         0.26, ">", theme, 18, "accent", True, "center")

    def render_timeline_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "阶段计划",
                         theme, slide_spec.get("kicker") or "TIMELINE")
        steps = trim_items(normalized_sections({"items": slide_spec.get(
            "steps") or slide_spec.get("items") or []}), 6, warnings, "timeline steps")
        add_rect(slide, 1.0, 3.65, 11.2, 0.04, theme_color(theme, "line"))
        gap = 11.0 / max(len(steps), 1)
        for index, step in enumerate(steps):
            left = 0.95 + index * gap
            add_rect(slide, left, 3.35, 0.42, 0.42,
                     theme_color(theme, "accent"), radius=True)
            phase = (step["raw"].get("phase") if isinstance(
                step.get("raw"), dict) else "") or f"阶段{index + 1}"
            add_text(slide, left - 0.35, 2.55, 1.25, 0.3, phase, theme, 9, "accent", True, "center",
                     density=density, overflow=overflow, warnings=warnings, context="timeline phase")
            add_text(slide, left - 0.55, 3.95, 1.65, 0.4, step["title"], theme, 11, "primary", True,
                     "center", density=density, overflow=overflow, warnings=warnings, context="timeline title")
            add_text(slide, left - 0.75, 4.48, 2.0, 0.95, "\n".join(step["points"][:2]), theme, 9, "secondary",
                     alignment="center", density=density, overflow=overflow, warnings=warnings, context="timeline detail")

    def render_metrics_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "关键指标",
                         theme, slide_spec.get("kicker") or "METRICS")
        metrics = trim_items(slide_spec.get("metrics") or slide_spec.get(
            "items") or [], 6, warnings, "metrics")
        for index, metric in enumerate(metrics):
            metric = metric if isinstance(metric, dict) else {
                "label": str(metric), "value": ""}
            row, col = divmod(index, 3)
            left, top = 0.85 + col * 4.15, 1.75 + row * 2.35
            add_rect(slide, left, top, 3.75, 1.85, theme_color(
                theme, "surface"), theme_color(theme, "line"), True)
            add_text(slide, left + 0.25, top + 0.22, 1.15, 0.3, metric.get("label") or metric.get("title") or "",
                     theme, 10, "muted", True, density=density, overflow=overflow, warnings=warnings, context="metric label")
            add_text(slide, left + 0.25, top + 0.58, 2.45, 0.55, metric.get("value") or metric.get("number") or "", theme, 24,
                     "accent", True, density=density, overflow=overflow, min_font_size=16, warnings=warnings, context="metric value")
            add_text(slide, left + 0.25, top + 1.22, 3.15, 0.42, metric.get("note") or metric.get("description") or "",
                     theme, 9, "secondary", density=density, overflow=overflow, warnings=warnings, context="metric note")

    def render_architecture_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "架构视图",
                         theme, slide_spec.get("kicker") or "ARCHITECTURE")
        layers = trim_items(slide_spec.get("layers") or slide_spec.get(
            "items") or [], 5, warnings, "architecture layers")
        layer_height = 4.9 / max(len(layers), 1)
        for index, layer in enumerate(layers):
            layer = layer if isinstance(layer, dict) else {
                "name": str(layer), "items": []}
            top = 1.55 + index * layer_height
            fill = theme_color(
                theme, "surface") if index % 2 == 0 else theme_color(theme, "light")
            add_rect(slide, 1.05, top, 11.1, layer_height - 0.18,
                     fill, theme_color(theme, "line"), True)
            add_text(slide, 1.32, top + 0.2, 1.7, 0.35, layer.get("name") or layer.get("title")
                     or f"Layer {index + 1}", theme, 12, "primary", True, density=density, overflow=overflow, warnings=warnings, context="architecture layer")
            items = trim_items(safe_lines(layer.get("items") or layer.get(
                "points") or layer.get("content"), 6), 5, warnings, "architecture layer items")
            for item_index, item in enumerate(items):
                item_left = 3.35 + item_index * 1.65
                add_rect(slide, item_left, top + 0.22, 1.35, 0.42, theme_color(theme, "primary")
                         if item_index % 2 == 0 else theme_color(theme, "secondary"), radius=True)
                add_text(slide, item_left + 0.08, top + 0.31, 1.18, 0.18, item, theme, 8, "surface", True, "center",
                         density=density, overflow=overflow, min_font_size=7, warnings=warnings, context="architecture item")

    def render_table_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "结构化信息",
                         theme, slide_spec.get("kicker") or "TABLE")
        table_spec = slide_spec.get("table") or {}
        headers = table_spec.get("headers") or []
        rows = table_spec.get("rows") or []
        if not headers and rows and isinstance(rows[0], dict):
            headers = list(rows[0].keys())
        rows = trim_items(rows, 6, warnings, "table rows")
        headers = trim_items(headers, 5, warnings, "table columns")
        col_width = 11.6 / max(len(headers), 1)
        row_height = 0.58
        top = 1.72
        for col, header in enumerate(headers):
            add_rect(slide, 0.85 + col * col_width, top, col_width, row_height,
                     theme_color(theme, "primary"), theme_color(theme, "line"))
            add_text(slide, 0.94 + col * col_width, top + 0.14, col_width - 0.18, 0.22, header, theme, 9, "surface",
                     True, "center", density=density, overflow=overflow, warnings=warnings, context="table header")
        for row_index, row in enumerate(rows):
            for col, header in enumerate(headers):
                value = row.get(header, "") if isinstance(row, dict) else (
                    row[col] if isinstance(row, list) and col < len(row) else "")
                y = top + row_height * (row_index + 1)
                add_rect(slide, 0.85 + col * col_width, y, col_width, row_height,
                         theme_color(theme, "surface"), theme_color(theme, "line"))
                add_text(slide, 0.94 + col * col_width, y + 0.13, col_width - 0.18, 0.24, value, theme, 8, "secondary",
                         alignment="center", density=density, overflow=overflow, min_font_size=7, warnings=warnings, context="table cell")

    def render_literature_matrix_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        table_spec = slide_spec.get("table") or {}
        rows = table_spec.get("rows") or slide_spec.get(
            "literature") or slide_spec.get("studies") or slide_spec.get("items") or []
        headers = table_spec.get("headers") or [
            "研究脉络", "主要观点", "不足之处", "本研究切入"]
        normalized_rows = []
        for row in rows:
            if isinstance(row, dict):
                normalized_rows.append([
                    row.get("theme") or row.get("topic") or row.get(
                        "direction") or row.get("title") or "",
                    row.get("argument") or row.get("view") or row.get(
                        "finding") or row.get("main_point") or row.get("观点") or "",
                    row.get("gap") or row.get(
                        "limitation") or row.get("不足") or "",
                    row.get("entry") or row.get("positioning") or row.get(
                        "contribution") or row.get("本研究切入") or "",
                ])
            else:
                normalized_rows.append(safe_lines(row, len(headers)))
        render_table_slide(
            presentation,
            {
                "title": slide_spec.get("title") or "文献综述：研究脉络与不足",
                "kicker": slide_spec.get("kicker") or "LITERATURE",
                "table": {"headers": headers, "rows": normalized_rows},
            },
            theme,
            "compact" if density == "standard" else density,
            overflow,
            warnings,
        )

    def render_research_questions_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "研究问题",
                         theme, slide_spec.get("kicker") or "RESEARCH QUESTIONS")
        background = slide_spec.get("background") or slide_spec.get("statement") or ""
        context = slide_spec.get("context") or ""
        if background:
            add_text(slide, 0.85, 1.38, 11.65, 0.34, background, theme, 10, "secondary", density=density,
                      overflow=overflow, warnings=warnings, context="research question background")
        if context:
            add_text(slide, 0.85, 1.76, 11.65, 0.42, context, theme, 9, "muted", density=density,
                     overflow=overflow, min_font_size=8, warnings=warnings, context="research question context")
        questions = trim_items(safe_lines(slide_spec.get("questions") or slide_spec.get(
            "research_questions") or slide_spec.get("items"), 6), 4, warnings, "research questions")
        top = 2.36 if context else 2.05
        question_gap = 0.86 if context else 1.12
        for index, question in enumerate(questions):
            y = top + index * question_gap
            add_rect(slide, 0.9, y, 0.58, 0.58, theme_color(
                theme, "accent"), radius=True)
            add_text(slide, 1.02, y + 0.15, 0.34, 0.22,
                     f"Q{index + 1}", theme, 10, "surface", True, "center")
            add_rect(slide, 1.65, y, 10.55, 0.62, theme_color(
                theme, "surface"), theme_color(theme, "line"), True)
            add_text(slide, 1.9, y + 0.14, 9.95, 0.28, question, theme, 13, "primary", True,
                     density=density, overflow=overflow, warnings=warnings, context="research question")
        gap = slide_spec.get("gap") or slide_spec.get("research_gap") or ""
        if gap:
            add_rect(slide, 0.9, 6.03, 11.3, 0.72, theme_color(
                theme, "light"), theme_color(theme, "line"), True)
            add_text(slide, 1.12, 6.18, 1.0, 0.24, "研究缺口", theme, 10, "accent", True,
                      density=density, overflow=overflow, warnings=warnings, context="research gap label")
            add_text(slide, 2.25, 6.12, 9.62, 0.42, gap, theme, 9, "secondary",
                     density=density, overflow=overflow, min_font_size=8, warnings=warnings, context="research gap")

    def render_theoretical_framework_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "理论框架",
                         theme, slide_spec.get("kicker") or "FRAMEWORK")
        framework_text = slide_spec.get("framework") or ""
        concepts = trim_items(normalized_sections({"items": slide_spec.get("concepts") or slide_spec.get(
            "variables") or slide_spec.get("items") or []}), 5, warnings, "framework concepts")
        if not concepts:
            concepts = [{"title": "核心概念", "points": safe_lines(
                framework_text or slide_spec.get("content"), 4), "raw": {}}]
        gap = 10.8 / max(len(concepts), 1)
        y = 1.82
        for index, concept in enumerate(concepts):
            left = 0.95 + index * gap
            width = min(2.0, gap - 0.18)
            add_rect(slide, left, y, width, 1.55, theme_color(
                theme, "surface"), theme_color(theme, "line"), True)
            add_rect(slide, left, y, width, 0.14, theme_color(theme, "accent"))
            add_text(slide, left + 0.12, y + 0.24, width - 0.24, 0.44, concept["title"], theme, 10, "primary",
                     True, "center", density=density, overflow=overflow, min_font_size=8, warnings=warnings, context="framework concept")
            add_text(slide, left + 0.12, y + 0.78, width - 0.24, 0.55, "\n".join(concept["points"][:2]), theme, 8, "secondary",
                      alignment="center", density=density, overflow=overflow, min_font_size=7, warnings=warnings, context="framework concept detail")
            if index < len(concepts) - 1:
                add_text(slide, left + width + 0.1, y + 0.43, 0.35,
                         0.28, ">", theme, 18, "accent", True, "center")
        propositions = safe_lines(slide_spec.get("relations") or slide_spec.get(
            "propositions") or slide_spec.get("hypotheses"), 4)
        add_card(slide, 0.95, 4.35, 5.45, 1.78, "关系假设 / 分析命题",
                 propositions, theme, "accent", density, overflow, warnings)
        mechanism_source = slide_spec.get("mechanism") or slide_spec.get(
            "explanation") or framework_text
        mechanism_title = "机制解释" if (slide_spec.get("mechanism") or slide_spec.get(
            "explanation")) else "总体框架"
        mechanism = safe_lines(mechanism_source, 4)
        add_card(slide, 6.75, 4.35, 5.45, 1.78, mechanism_title, mechanism,
                 theme, "success", density, overflow, warnings)

    def render_method_design_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "研究设计",
                         theme, slide_spec.get("kicker") or "METHOD")
        method_items = []
        if slide_spec.get("data_sources") or slide_spec.get("data"):
            method_items.append({"title": "数据来源", "points": safe_lines(
                slide_spec.get("data_sources") or slide_spec.get("data"), 3)})
        if slide_spec.get("sample") or slide_spec.get("scope"):
            method_items.append({"title": "样本范围", "points": safe_lines(
                slide_spec.get("sample") or slide_spec.get("scope"), 3)})
        if slide_spec.get("variables"):
            method_items.append({"title": "变量设计", "points": safe_lines(
                slide_spec.get("variables"), 3)})

        explicit_methods = slide_spec.get("methods") or slide_spec.get(
            "research_design") or slide_spec.get("steps")
        if explicit_methods:
            if isinstance(explicit_methods, list) and all(isinstance(item, dict) for item in explicit_methods):
                method_items.extend(explicit_methods)
            else:
                method_items.append({"title": "分析方法", "points": safe_lines(
                    explicit_methods, 3)})
        elif slide_spec.get("analysis") or slide_spec.get("method"):
            method_items.append({"title": "分析方法", "points": safe_lines(
                slide_spec.get("analysis") or slide_spec.get("method"), 3)})

        if not method_items:
            method_items = [
                {"title": "数据来源", "points": safe_lines(slide_spec.get(
                    "data_sources") or slide_spec.get("data"), 3)},
                {"title": "样本范围", "points": safe_lines(
                    slide_spec.get("sample") or slide_spec.get("scope"), 3)},
                {"title": "变量设计", "points": safe_lines(
                    slide_spec.get("variables"), 3)},
                {"title": "分析方法", "points": safe_lines(
                    slide_spec.get("analysis") or slide_spec.get("method"), 3)},
            ]
        steps = trim_items(normalized_sections(
            {"items": method_items}), 4, warnings, "method design steps")
        positions = [(0.85, 1.72), (6.85, 1.72), (0.85, 4.05), (6.85, 4.05)]
        labels = ["01", "02", "03", "04"]
        for index, step in enumerate(steps):
            left, top = positions[index]
            add_rect(slide, left, top, 5.2, 1.78, theme_color(
                theme, "surface"), theme_color(theme, "line"), True)
            add_rect(slide, left + 0.22, top + 0.22, 0.52, 0.52,
                     theme_color(theme, "primary"), radius=True)
            add_text(slide, left + 0.31, top + 0.36, 0.34, 0.18,
                     labels[index], theme, 9, "surface", True, "center")
            add_text(slide, left + 0.9, top + 0.24, 3.95, 0.32, step["title"], theme, 13, "primary",
                     True, density=density, overflow=overflow, warnings=warnings, context="method step title")
            add_text(slide, left + 0.9, top + 0.78, 3.95, 0.68, "\n".join([f"- {line}" for line in step["points"][:3]]),
                     theme, 9, "secondary", density=density, overflow=overflow, warnings=warnings, context="method step body")

    def render_findings_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "主要发现",
                         theme, slide_spec.get("kicker") or "FINDINGS")
        findings = trim_items(normalized_sections({"items": slide_spec.get("findings") or slide_spec.get(
            "items") or slide_spec.get("sections") or []}), 4, warnings, "findings")
        headline = slide_spec.get(
            "headline") or slide_spec.get("statement") or ""
        if headline:
            add_rect(slide, 0.9, 1.55, 11.25, 0.88, theme_color(
                theme, "light"), theme_color(theme, "line"), True)
            add_text(slide, 1.15, 1.78, 10.75, 0.34, headline, theme, 16, "primary", True,
                     density=density, overflow=overflow, warnings=warnings, context="finding headline")
        positions = [(0.9, 2.75), (6.85, 2.75), (0.9, 4.78), (6.85, 4.78)] if headline else [
            (0.9, 1.72), (6.85, 1.72), (0.9, 4.05), (6.85, 4.05)]
        card_height = 1.42 if headline else 1.78
        for index, finding in enumerate(findings):
            left, top = positions[index]
            add_card(slide, left, top, 5.25, card_height, finding["title"], finding["points"],
                     theme, "accent" if index % 2 == 0 else "success", density, overflow, warnings)

    def render_contribution_limitations_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_slide_header(slide, slide_spec.get("title") or "研究贡献与局限",
                         theme, slide_spec.get("kicker") or "CONTRIBUTION")
        contributions = safe_lines(slide_spec.get("contributions") or slide_spec.get(
            "contribution") or slide_spec.get("left") or slide_spec.get("items"), 6)
        limitations = safe_lines(slide_spec.get("limitations") or slide_spec.get(
            "limitation") or slide_spec.get("right"), 6)
        implications = safe_lines(slide_spec.get("implications") or slide_spec.get(
            "future") or slide_spec.get("outlook"), 4)
        add_card(slide, 0.85, 1.72, 5.45, 3.1, "研究贡献", contributions,
                 theme, "success", density, overflow, warnings)
        add_card(slide, 6.85, 1.72, 5.45, 3.1, "局限与边界", limitations,
                 theme, "danger", density, overflow, warnings)
        if implications:
            add_card(slide, 0.85, 5.28, 11.45, 1.05, "后续研究方向",
                     implications, theme, "accent", density, overflow, warnings)

    def render_quote_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_theme_background(slide, theme)
        add_rect(slide, 1.0, 1.45, 0.12, 4.55, theme_color(theme, "accent"))
        add_text(slide, 1.35, 1.35, 10.2, 0.45, slide_spec.get("title") or "核心结论", theme, 18, "primary",
                 True, density=density, overflow=overflow, warnings=warnings, context="quote title")
        statement = slide_spec.get(
            "statement") or slide_spec.get("content") or ""
        add_text(slide, 1.35, 2.25, 10.2, 1.8, statement, theme, 24, "primary", True, density=density,
                 overflow=overflow, min_font_size=15, warnings=warnings, context="quote statement")
        points = safe_lines(slide_spec.get("points")
                            or merge_point_lines(slide_spec, slide_spec.get("items")), 4)
        if points:
            add_text(slide, 1.42, 4.62, 9.9, 0.75, "  /  ".join(points), theme, 11, "secondary",
                     density=density, overflow=overflow, warnings=warnings, context="quote points")

    def render_closing_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_theme_background(slide, theme)
        title = slide_spec.get("title") or "谢谢"
        subtitle = slide_spec.get(
            "subtitle") or slide_spec.get("content") or ""
        add_text(slide, 2.0, 2.35, 9.3, 0.8, title, theme, 32, "primary", True, "center", density=density,
                 overflow=overflow, min_font_size=18, warnings=warnings, context="closing title")
        add_text(slide, 2.35, 3.25, 8.6, 0.7, subtitle, theme, 14, "secondary", alignment="center",
                 density=density, overflow=overflow, warnings=warnings, context="closing subtitle")

    def add_expert_header(
        slide,
        title: str,
        theme: Dict[str, Any],
        kicker: str = "",
        density: str = "standard",
        overflow: str = "shrink_then_truncate",
        warnings: Optional[List[str]] = None,
    ) -> None:
        add_rect(slide, 0, 0, 13.333, 7.5, theme_color(theme, "background"))
        add_rect(slide, 0, 0, 13.333, 0.16, theme_color(theme, "primary"))
        add_rect(slide, 0, 0.16, 0.16, 7.34, theme_color(theme, "accent"))
        if kicker:
            add_text(slide, 0.72, 0.34, 3.0, 0.22, kicker,
                     theme, 8, "accent", True, density=density, overflow=overflow, warnings=warnings)
        add_text(slide, 2.2, 0.48, 8.9, 0.55, title, theme, fit_text_size(
            title, 22, 16, 26), "primary", True, "center", density=density, overflow=overflow, warnings=warnings)
        add_rect(slide, 5.88, 1.18, 1.55, 0.04, theme_color(theme, "accent"))

    def render_expert_section_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_rect(slide, 0, 0, 13.333, 7.5, theme_color(theme, "background"))
        add_rect(slide, 0, 0, 13.333, 0.18, theme_color(theme, "primary"))
        add_rect(slide, 0, 7.28, 13.333, 0.22, theme_color(theme, "primary"))
        add_rect(slide, 0.72, 1.25, 1.55, 0.08, theme_color(theme, "accent"))
        title = slide_spec.get("title") or "章节标题"
        subtitle = slide_spec.get("subtitle") or slide_spec.get("content") or ""
        add_text(slide, 2.95, 2.08, 7.45, 0.92, title, theme, fit_text_size(
            title, 30, 20, 24), "primary", True, "center", density=density, overflow=overflow, warnings=warnings, context="expert section title")
        if subtitle:
            add_text(slide, 2.9, 3.05, 7.55, 0.45, subtitle, theme, 13, "secondary",
                     alignment="center", density=density, overflow=overflow, warnings=warnings, context="expert section subtitle")
        tag = slide_spec.get("tag")
        if tag:
            add_text(slide, 0.0, 5.05, 13.333, 0.36, tag, theme, 9, "muted",
                     alignment="center", density=density, overflow=overflow, warnings=warnings, context="expert section tag")

    def render_expert_title_content_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_expert_header(slide, slide_spec.get("title") or "内容页",
                          theme, density=density, overflow=overflow, warnings=warnings)
        statement = slide_spec.get("statement") or slide_spec.get("content") or ""
        if statement:
            add_text(slide, 1.1, 1.5, 11.1, 0.55, statement, theme, 16, "primary", True,
                     alignment="center", density=density, overflow=overflow, warnings=warnings, context="expert content statement")
        sections = trim_items(normalized_sections(
            {"items": slide_spec.get("items") or slide_spec.get("sections") or []}), 4, warnings, "expert content blocks")
        if not sections and not statement:
            sections = [{"title": "核心内容", "points": safe_lines(slide_spec.get("text"), 4), "raw": {}}]
        positions = [(1.0, 2.35), (7.05, 2.35), (1.0, 4.55), (7.05, 4.55)]
        for index, section in enumerate(sections):
            left, top = positions[index]
            add_rect(slide, left, top, 5.3, 1.5, theme_color(theme, "surface"),
                     theme_color(theme, "line"), True)
            add_rect(slide, left, top, 0.08, 1.5, theme_color(
                theme, "accent" if index % 2 == 0 else "primary"))
            add_text(slide, left + 0.22, top + 0.18, 4.8, 0.3, section["title"], theme, 13, "primary",
                     True, density=density, overflow=overflow, warnings=warnings, context="expert content block title")
            add_text(slide, left + 0.22, top + 0.62, 4.85, 0.7, "\n".join([f"- {line}" for line in section["points"][:3]]),
                     theme, 9, "secondary", density=density, overflow=overflow, warnings=warnings, context="expert content block body")

    def render_expert_split_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_expert_header(slide, slide_spec.get("title") or "图文页",
                          theme, density=density, overflow=overflow, warnings=warnings)
        image_path = slide_spec.get("image_path", "")
        if image_path and os.path.exists(image_path):
            ppt_utils.add_image(slide, image_path, 0.95, 1.55, 5.2, 4.7)
        else:
            add_rect(slide, 0.95, 1.55, 5.2, 4.7, theme_color(
                theme, "surface"), theme_color(theme, "line"), True)
            add_text(slide, 1.25, 3.68, 4.6, 0.35, slide_spec.get("image_caption") or "图示区域",
                     theme, 12, "muted", True, "center", density=density, overflow=overflow, warnings=warnings, context="expert split placeholder")
        if slide_spec.get("image_caption") and not image_path:
            add_text(slide, 1.15, 6.35, 4.8, 0.25, slide_spec.get("image_caption"),
                     theme, 8, "muted", alignment="center", density=density, overflow=overflow)
        sections = trim_items(normalized_sections(
            {"items": slide_spec.get("items") or slide_spec.get("sections") or []}), 4, warnings, "expert split blocks")
        for index, section in enumerate(sections):
            top = 1.55 + index * 1.18
            add_text(slide, 6.75, top, 4.9, 0.3, section["title"], theme, 13, "primary", True,
                     density=density, overflow=overflow, warnings=warnings, context="expert split title")
            add_text(slide, 6.75, top + 0.38, 5.05, 0.55, "\n".join([f"- {line}" for line in section["points"][:2]]),
                     theme, 9, "secondary", density=density, overflow=overflow, warnings=warnings, context="expert split body")
            add_rect(slide, 6.48, top + 0.06, 0.08, 0.72, theme_color(theme, "accent"))

    def render_expert_path_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_expert_header(slide, slide_spec.get("title") or "流程路径",
                          theme, density=density, overflow=overflow, warnings=warnings)
        steps = trim_items(normalized_sections(
            {"items": slide_spec.get("steps") or slide_spec.get("items") or slide_spec.get("sections") or []}), 6, warnings, "expert path steps")
        for index, step in enumerate(steps):
            y = 1.55 + index * 0.82
            add_rect(slide, 5.94, y, 0.48, 0.48, theme_color(theme, "accent"), radius=True)
            add_text(slide, 6.05, y + 0.12, 0.25, 0.2, str(index + 1), theme,
                     9, "surface", True, "center", density=density, overflow=overflow)
            add_text(slide, 1.05, y + 0.05, 4.25, 0.26, step["title"], theme, 12, "primary", True,
                     density=density, overflow=overflow, warnings=warnings, context="expert path title")
            detail = "；".join(step["points"][:2])
            add_text(slide, 6.85, y + 0.02, 5.05, 0.35, detail, theme, 9, "secondary",
                     density=density, overflow=overflow, warnings=warnings, context="expert path detail")

    def render_expert_scope_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> None:
        slide = make_blank_slide(presentation)
        add_expert_header(slide, slide_spec.get("title") or "范围清单",
                          theme, density=density, overflow=overflow, warnings=warnings)
        items = trim_items(normalized_sections(
            {"items": slide_spec.get("items") or slide_spec.get("sections") or []}), 4, warnings, "expert scope items")
        for index, item in enumerate(items):
            top = 1.55 + index * 1.2
            add_rect(slide, 1.25, top, 10.75, 0.9, theme_color(
                theme, "surface"), theme_color(theme, "line"), True)
            add_rect(slide, 1.25, top, 0.12, 0.9, theme_color(theme, "accent"))
            add_text(slide, 1.62, top + 0.18, 2.7, 0.26, item["title"], theme, 13, "primary", True,
                     density=density, overflow=overflow, warnings=warnings, context="expert scope title")
            add_text(slide, 4.55, top + 0.15, 6.9, 0.42, "；".join(item["points"][:2]), theme, 10, "secondary",
                     density=density, overflow=overflow, warnings=warnings, context="expert scope detail")

    def render_generated_slide(
        presentation,
        slide_spec: Dict[str, Any],
        theme: Dict[str, Any],
        density: str,
        overflow: str,
        warnings: List[str],
    ) -> str:
        slide_type = infer_slide_type(slide_spec)
        if slide_type == "cover":
            render_cover_slide(presentation, slide_spec,
                               theme, density, overflow, warnings)
        elif slide_type == "expert_section":
            render_expert_section_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type == "expert_title_content":
            render_expert_title_content_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type == "expert_split":
            render_expert_split_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type == "expert_path":
            render_expert_path_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type == "expert_scope":
            render_expert_scope_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type in {"comparison"}:
            render_two_column_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type in {"process"}:
            render_process_slide(presentation, slide_spec,
                                 theme, density, overflow, warnings)
        elif slide_type == "timeline":
            render_timeline_slide(presentation, slide_spec,
                                  theme, density, overflow, warnings)
        elif slide_type in {"cards", "section"}:
            render_cards_slide(presentation, slide_spec,
                               theme, density, overflow, warnings)
        elif slide_type == "metrics":
            render_metrics_slide(presentation, slide_spec,
                                 theme, density, overflow, warnings)
        elif slide_type == "architecture":
            render_architecture_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type == "table":
            render_table_slide(presentation, slide_spec,
                               theme, density, overflow, warnings)
        elif slide_type == "literature_matrix":
            render_literature_matrix_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type == "research_questions":
            render_research_questions_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type == "theoretical_framework":
            render_theoretical_framework_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type == "method_design":
            render_method_design_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type == "findings":
            render_findings_slide(presentation, slide_spec,
                                  theme, density, overflow, warnings)
        elif slide_type == "contribution_limitations":
            render_contribution_limitations_slide(
                presentation, slide_spec, theme, density, overflow, warnings)
        elif slide_type == "quote":
            render_quote_slide(presentation, slide_spec,
                               theme, density, overflow, warnings)
        elif slide_type in {"closing", "end"}:
            render_closing_slide(presentation, slide_spec,
                                 theme, density, overflow, warnings)
        else:
            render_summary_slide(presentation, slide_spec,
                                 theme, density, overflow, warnings)
            slide_type = "summary"
        if presentation.slides:
            add_source_note(presentation.slides[-1], slide_spec, theme, warnings)
        return slide_type

    def collect_template_files(template_directory: str = "") -> List[str]:
        directories = []
        if template_directory:
            candidate = os.path.abspath(os.path.expanduser(template_directory))
            if os.path.isdir(candidate):
                directories.append(candidate)
        else:
            for directory in get_template_search_directories():
                candidate = os.path.abspath(os.path.expanduser(directory))
                if os.path.isdir(candidate) and candidate not in directories:
                    directories.append(candidate)

        template_files: List[str] = []
        for directory in directories:
            for root, _, files in os.walk(directory):
                for file_name in files:
                    if file_name.lower().endswith((".pptx", ".potx")):
                        template_files.append(os.path.abspath(
                            os.path.join(root, file_name)))

        return sorted(set(template_files))

    def summarize_layout(layout, layout_index: int, include_placeholders: bool = False) -> Dict[str, Any]:
        layout_info: Dict[str, Any] = {
            "layout_index": layout_index,
            "layout_name": getattr(layout, "name", f"Layout {layout_index}"),
            "placeholder_count": len(layout.placeholders),
        }

        if include_placeholders:
            placeholders = []
            for placeholder in layout.placeholders:
                placeholders.append({
                    "idx": getattr(placeholder.placeholder_format, "idx", None),
                    "type": str(getattr(placeholder.placeholder_format, "type", "")),
                    "name": getattr(placeholder, "name", ""),
                })
            layout_info["placeholders"] = placeholders

        layout_info["recommended_for"] = recommend_slide_types_for_layout(
            layout_info["layout_name"],
            layout_info["placeholder_count"],
        )
        return layout_info

    def recommend_slide_types_for_layout(layout_name: str, placeholder_count: int) -> List[str]:
        normalized = normalize_text(layout_name)
        recommendations: List[str] = []

        if "标题幻灯片" in layout_name or "封面" in layout_name:
            recommendations.append("cover")
        if "title" in normalized and ("slide" in normalized or "cover" in normalized):
            recommendations.append("cover")
        if "section" in normalized:
            recommendations.append("section")
        if "two" in normalized or "comparison" in normalized:
            recommendations.append("two_column")
        if "picture" in normalized or "image" in normalized:
            recommendations.extend(["image_text", "cover"])
        if "chart" in normalized:
            recommendations.append("chart")
        if "table" in normalized:
            recommendations.append("table")
        if ("content" in normalized or "text" in normalized) and "two" not in normalized:
            recommendations.extend(["bullet", "agenda", "summary"])

        if not recommendations:
            if placeholder_count >= 3:
                recommendations.extend(["bullet", "two_column"])
            elif placeholder_count >= 2:
                recommendations.extend(["bullet", "summary"])
            else:
                recommendations.append("section")

        deduped: List[str] = []
        for item in recommendations:
            if item not in deduped:
                deduped.append(item)
        return deduped

    def summarize_template_file(template_path: str, include_layouts: bool, include_placeholders: bool) -> Dict[str, Any]:
        info = ppt_utils.get_template_info(template_path)
        template_name = os.path.basename(template_path)

        template_summary: Dict[str, Any] = {
            "template_id": slugify(os.path.splitext(template_name)[0]),
            "template_name": template_name,
            "template_path": template_path,
            "file_size_bytes": info["file_size_bytes"],
            "layout_count": info["layout_count"],
            "slide_count": info["slide_count"],
            "core_properties": info["core_properties"],
            "style_tags": infer_style_tags(template_name, info["core_properties"]),
        }

        if include_layouts:
            presentation = ppt_utils.open_presentation(template_path)
            template_summary["layouts"] = [
                summarize_layout(layout, idx, include_placeholders)
                for idx, layout in enumerate(presentation.slide_layouts)
            ]

        return template_summary

    def infer_style_tags(template_name: str, core_properties: Dict[str, Any]) -> List[str]:
        text = " ".join([
            template_name,
            core_properties.get("title") or "",
            core_properties.get("subject") or "",
            core_properties.get("keywords") or "",
        ]).lower()

        rules = {
            "tech": ["tech", "technology", "科技", "数字化", "蓝"],
            "business": ["business", "商务", "report", "proposal", "方案"],
            "education": ["education", "teaching", "教学", "academic", "研究"],
            "party": ["党建", "government", "党"],
            "dark": ["dark", "深蓝", "black"],
            "blue": ["blue", "蓝"],
        }

        tags = [tag for tag, hints in rules.items() if any(
            hint in text for hint in hints)]
        return tags or ["general"]

    def resolve_template_reference(template_config: Optional[Dict[str, Any]] = None) -> Tuple[Optional[str], Dict[str, Any]]:
        template_config = template_config or {}
        template_id = (template_config.get("template_id") or "").strip()
        template_path = (template_config.get("template_path") or "").strip()
        template_name = (template_config.get("template_name") or "").strip()
        theme_hint = (template_config.get("theme_hint") or "").strip()
        template_directory = (template_config.get(
            "template_directory") or "").strip()

        if template_path:
            absolute_path = os.path.abspath(os.path.expanduser(template_path))
            if os.path.isfile(absolute_path):
                return absolute_path, {
                    "selection_method": "template_path",
                    "mode": "strong_template",
                }
            return None, {
                "error": f"Template file not found: {template_path}",
                "selection_method": "template_path",
            }

        template_files = collect_template_files(template_directory)
        if not template_files:
            return None, {
                "selection_method": "fallback",
                "mode": "blank_presentation",
            }

        normalized_id = normalize_text(template_id)
        normalized_name = normalize_text(template_name)
        normalized_theme = normalize_text(theme_hint)

        best_match: Optional[str] = None
        best_reason = ""

        for candidate in template_files:
            file_name = os.path.basename(candidate)
            stem = os.path.splitext(file_name)[0]
            candidate_id = slugify(stem)

            if normalized_id and normalize_text(candidate_id) == normalized_id:
                best_match = candidate
                best_reason = "template_id"
                break

            if normalized_name and normalize_text(file_name) == normalized_name:
                best_match = candidate
                best_reason = "template_name"
                break

        if not best_match and normalized_theme:
            scored_candidates: List[Tuple[int, str]] = []
            for candidate in template_files:
                file_name = os.path.basename(candidate)
                score = 0
                normalized_file_name = normalize_text(file_name)
                if normalized_theme in normalized_file_name:
                    score += 4
                for token in re.split(r"[\s,_\-]+", theme_hint.lower()):
                    token = token.strip()
                    if token and token in file_name.lower():
                        score += 1
                if score:
                    scored_candidates.append((score, candidate))
            if scored_candidates:
                scored_candidates.sort(key=lambda item: (-item[0], item[1]))
                best_match = scored_candidates[0][1]
                best_reason = "theme_hint"

        if not best_match:
            best_match = template_files[0]
            best_reason = "default_template"

        return best_match, {
            "selection_method": best_reason,
            "mode": "strong_template",
        }

    def build_layout_catalog(presentation) -> List[Dict[str, Any]]:
        return [
            summarize_layout(layout, idx, include_placeholders=True)
            for idx, layout in enumerate(presentation.slide_layouts)
        ]

    def placeholder_role(placeholder) -> str:
        placeholder_type = str(
            getattr(placeholder.placeholder_format, "type", "")).upper()
        placeholder_name = normalize_text(getattr(placeholder, "name", ""))

        if "SUBTITLE" in placeholder_type or "subtitle" in placeholder_name:
            return "subtitle"
        if "TITLE" in placeholder_type or "title" in placeholder_name:
            return "title"
        if any(token in placeholder_type for token in ("BODY", "CONTENT", "OBJECT")):
            return "body"
        if any(token in placeholder_type for token in ("DATE", "FOOTER", "SLIDE_NUMBER")):
            return "system"
        return "text"

    def layout_placeholder_map(layout) -> Dict[str, Any]:
        placeholders = []
        for placeholder in layout.placeholders:
            placeholders.append({
                "idx": getattr(placeholder.placeholder_format, "idx", None),
                "type": str(getattr(placeholder.placeholder_format, "type", "")),
                "name": getattr(placeholder, "name", ""),
                "role": placeholder_role(placeholder),
            })
        return {
            "placeholders": placeholders,
            "roles": sorted({item["role"] for item in placeholders}),
        }

    def select_layout_index(
        presentation,
        slide_type: str,
        preferred_layout: Optional[Dict[str, Any]] = None,
        policy: str = "preferred_then_best_match",
    ) -> int:
        layouts = build_layout_catalog(presentation)

        if preferred_layout and preferred_layout.get("layout_index") is not None:
            preferred_index = preferred_layout["layout_index"]
            if 0 <= preferred_index < len(layouts):
                return preferred_index
            if policy == "strict_preferred_only":
                raise ValueError(
                    f"Preferred layout index is invalid: {preferred_index}")

        normalized_slide_type = normalize_text(slide_type)

        for layout in layouts:
            if slide_type in layout["recommended_for"] or normalized_slide_type in [
                normalize_text(item) for item in layout["recommended_for"]
            ]:
                return layout["layout_index"]

        for layout in layouts:
            name = normalize_text(layout["layout_name"])
            if normalized_slide_type in name:
                return layout["layout_index"]

        if policy == "strict_preferred_only":
            raise ValueError(f"No layout matched slide type '{slide_type}'")

        return 0 if slide_type in ("cover", "section") else min(1, len(layouts) - 1)

    def clear_all_slides(presentation) -> int:
        removed = len(presentation.slides)
        for slide_id in list(presentation.slides._sldIdLst):
            relationship_id = slide_id.rId
            presentation.part.drop_rel(relationship_id)
            presentation.slides._sldIdLst.remove(slide_id)
        return removed

    def remove_slides_from(presentation, start_index: int) -> int:
        removed = 0
        while len(presentation.slides) > start_index:
            slide_id = presentation.slides._sldIdLst[start_index]
            relationship_id = slide_id.rId
            presentation.part.drop_rel(relationship_id)
            presentation.slides._sldIdLst.remove(slide_id)
            removed += 1
        return removed

    def is_prompt_text(text: str) -> bool:
        normalized = (text or "").strip().lower()
        if not normalized:
            return False
        prompts = [
            "click to add",
            "单击此处添加",
            "点击此处添加",
            "add title",
            "add subtitle",
        ]
        return any(prompt in normalized for prompt in prompts)

    def clear_text_frame(text_frame) -> None:
        text_frame.clear()
        if text_frame.paragraphs:
            text_frame.paragraphs[0].text = ""

    def clear_slide_prompt_text(slide) -> None:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            if is_prompt_text(getattr(shape.text_frame, "text", "") or ""):
                clear_text_frame(shape.text_frame)

    def reset_slide_text_content(slide) -> None:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            if hasattr(shape, "is_placeholder") and shape.is_placeholder and placeholder_role(shape) == "system":
                continue
            clear_text_frame(shape.text_frame)

    def get_text_placeholders(slide) -> Dict[str, List[Any]]:
        placeholders: Dict[str, List[Any]] = {
            "title": [], "subtitle": [], "body": [], "text": []}
        for placeholder in slide.placeholders:
            if not hasattr(placeholder, "text_frame"):
                continue
            role = placeholder_role(placeholder)
            if role in placeholders:
                placeholders[role].append(placeholder)
            elif role != "system":
                placeholders["text"].append(placeholder)

        generic_text_shapes = []
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                continue
            existing_text = (shape.text_frame.text or "").strip()
            shape_type_name = str(getattr(shape, "shape_type", ""))
            if not existing_text and "TEXT_BOX" not in shape_type_name:
                continue
            generic_text_shapes.append(shape)

        generic_text_shapes.sort(key=lambda item: (
            item.top, -(item.width * item.height)))
        for shape in generic_text_shapes:
            if shape.top < int(2.2 * emu_per_inch) and not placeholders["title"]:
                placeholders["title"].append(shape)
            else:
                placeholders["text"].append(shape)
        return placeholders

    def set_shape_text(shape, text: str) -> bool:
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            return False
        clear_text_frame(shape.text_frame)
        if text:
            shape.text = text
        return True

    def apply_field_format(shape, field_spec: Dict[str, Any]) -> None:
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            return

        text_format = field_spec.get("format", {})
        font_size = field_spec.get("font_size", text_format.get("font_size"))
        font_name = field_spec.get("font_name", text_format.get("font_name"))
        bold = field_spec.get("bold", text_format.get("bold"))
        italic = field_spec.get("italic", text_format.get("italic"))
        color = field_spec.get("color", text_format.get("color"))
        alignment = field_spec.get("alignment", text_format.get("alignment"))

        if any(value is not None for value in [font_size, font_name, bold, italic, color, alignment]):
            ppt_utils.format_text_advanced(
                shape.text_frame,
                font_size=font_size,
                font_name=font_name,
                bold=bold,
                italic=italic,
                color=tuple(color) if isinstance(
                    color, list) and len(color) == 3 else None,
                alignment=alignment,
            )

    def join_slide_items(slide_spec: Dict[str, Any]) -> List[str]:
        items = slide_spec.get("items") or slide_spec.get(
            "bullet_points") or slide_spec.get("points") or []
        if not isinstance(items, list):
            items = [str(items)]
        return [str(item) for item in items if str(item).strip()]

    def render_known_template_slide(slide, slide_spec: Dict[str, Any]) -> bool:
        slide_type = (slide_spec.get("slide_type")
                      or slide_spec.get("type") or "").strip()
        layout_name = getattr(slide.slide_layout, "name", "")
        shape_count = len(slide.shapes)
        items = join_slide_items(slide_spec)
        title = slide_spec.get("title", "")
        subtitle = slide_spec.get("subtitle", "")

        if layout_name == "标题幻灯片" and shape_count >= 15 and slide_type in ("cover", "closing"):
            updated = False
            updated |= set_shape_text(slide.shapes[11], title)
            updated |= set_shape_text(slide.shapes[14], subtitle)
            return updated

        if layout_name == "标题和内容" and shape_count == 15 and slide_type == "agenda":
            updated = False
            updated |= set_shape_text(slide.shapes[5], title or "目录")
            agenda_slots = [7, 9, 11, 13]
            for idx, item in zip(agenda_slots, items[:4]):
                updated |= set_shape_text(slide.shapes[idx], item)
            for idx in agenda_slots[len(items[:4]):]:
                updated |= set_shape_text(slide.shapes[idx], "")
            return updated

        if layout_name == "标题和内容" and shape_count == 8 and slide_type in ("section", "summary"):
            updated = False
            updated |= set_shape_text(slide.shapes[2], title)
            body_text = slide_spec.get("content") or slide_spec.get(
                "text") or "；".join(items[:3])
            updated |= set_shape_text(slide.shapes[4], body_text)
            return updated

        if layout_name == "Main" and shape_count >= 29 and slide_type in ("bullet", "summary"):
            updated = False
            updated |= set_shape_text(slide.shapes[25], title)
            slots = [19, 20, 21, 22]
            for idx, item in zip(slots, items[:4]):
                updated |= set_shape_text(slide.shapes[idx], item)
            for idx in slots[len(items[:4]):]:
                updated |= set_shape_text(slide.shapes[idx], "")
            if items:
                updated |= set_shape_text(slide.shapes[26], items[0])
            else:
                updated |= set_shape_text(slide.shapes[26], "")
            return updated

        return False

    def fill_title_placeholders(slide, title: str = "", subtitle: str = "") -> Dict[str, bool]:
        placeholders = get_text_placeholders(slide)
        title_filled = False
        subtitle_filled = False

        if title:
            title_placeholder = slide.shapes.title
            if title_placeholder and hasattr(title_placeholder, "text_frame"):
                clear_text_frame(title_placeholder.text_frame)
                title_placeholder.text = title
                title_filled = True
            elif placeholders["title"]:
                clear_text_frame(placeholders["title"][0].text_frame)
                placeholders["title"][0].text = title
                title_filled = True

        subtitle_targets = placeholders["subtitle"] or placeholders["body"][:1] or placeholders["text"][:1]
        if subtitle_targets:
            clear_text_frame(subtitle_targets[0].text_frame)
            if subtitle:
                subtitle_targets[0].text = subtitle
                subtitle_filled = True

        clear_slide_prompt_text(slide)
        return {"title_filled": title_filled, "subtitle_filled": subtitle_filled}

    def find_body_placeholders(slide) -> List[Any]:
        placeholders = get_text_placeholders(slide)
        title_placeholder = slide.shapes.title
        return [item for item in placeholders["body"] + placeholders["text"] if item != title_placeholder]

    def add_bullet_content(slide, lines: List[str], use_all_body_placeholders: bool = True) -> bool:
        body_placeholders = find_body_placeholders(slide)
        if not body_placeholders:
            return False

        if use_all_body_placeholders and len(body_placeholders) > 1 and len(lines) > 2:
            chunks: List[List[str]] = [[] for _ in body_placeholders]
            for index, line in enumerate(lines):
                chunks[index % len(body_placeholders)].append(str(line))

            for placeholder, chunk in zip(body_placeholders, chunks):
                text_frame = placeholder.text_frame
                clear_text_frame(text_frame)
                for item_index, line in enumerate(chunk):
                    paragraph = text_frame.paragraphs[0] if item_index == 0 else text_frame.add_paragraph(
                    )
                    paragraph.text = line
                    paragraph.level = 0
            return True

        text_frame = body_placeholders[0].text_frame
        clear_text_frame(text_frame)
        for index, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph(
            )
            paragraph.text = str(line)
            paragraph.level = 0
        return True

    def get_slide_for_render(
        presentation,
        slide_spec: Dict[str, Any],
        template_rendering: Dict[str, Any],
        slide_position: int,
    ):
        reuse_existing_template_slides = template_rendering.get(
            "use_template_sample_slides", True)
        existing_template_slide_count = template_rendering.get(
            "existing_template_slide_count", 0)

        if reuse_existing_template_slides and slide_position < existing_template_slide_count and slide_position < len(presentation.slides):
            slide = presentation.slides[slide_position]
            reset_slide_text_content(slide)
            return slide, slide.slide_layout, slide_position, True

        layout_index = select_layout_index(
            presentation,
            (slide_spec.get("slide_type") or slide_spec.get(
                "type") or "bullet").strip(),
            slide_spec.get("preferred_layout"),
            template_rendering.get(
                "layout_selection_policy", "preferred_then_best_match"),
        )
        slide, layout = ppt_utils.add_slide(presentation, layout_index)
        return slide, layout, len(presentation.slides) - 1, False

    def render_slide_from_spec(
        presentation,
        slide_spec: Dict[str, Any],
        template_rendering: Dict[str, Any],
        slide_position: int,
    ) -> Dict[str, Any]:
        slide_type = (slide_spec.get("slide_type")
                      or slide_spec.get("type") or "bullet").strip()
        placeholder_policy = template_rendering.get(
            "placeholder_fill_policy", "prefer_placeholders")
        slide, layout, slide_index, reused_template_slide = get_slide_for_render(
            presentation,
            slide_spec,
            template_rendering,
            slide_position,
        )
        layout_index = next((idx for idx, candidate in enumerate(
            presentation.slide_layouts) if candidate == layout), -1)

        title = slide_spec.get("title", "")
        subtitle = slide_spec.get("subtitle", "")

        warnings: List[str] = []
        rendered_using = "template_layout"
        known_template_rendered = render_known_template_slide(
            slide, slide_spec)
        if known_template_rendered:
            clear_slide_prompt_text(slide)
            return {
                "slide_index": slide_index,
                "slide_type": slide_type,
                "layout_index": layout_index,
                "layout_name": getattr(layout, "name", f"Layout {layout_index}"),
                "rendered_using": "template_layout",
                "reused_template_slide": reused_template_slide,
                "warnings": warnings,
            }

        fill_result = fill_title_placeholders(
            slide, title=title, subtitle=subtitle)

        if slide_type in ("cover", "section"):
            if title and not fill_result["title_filled"]:
                ppt_utils.add_textbox(
                    slide, 0.8, 0.8, 8.5, 1.0, title, font_size=24, bold=True)
                warnings.append(
                    "Cover title used dynamic textbox because no title placeholder was available")
                rendered_using = "dynamic_content"
            if subtitle and not fill_result["subtitle_filled"]:
                ppt_utils.add_textbox(
                    slide, 0.9, 1.8, 8.0, 0.8, subtitle, font_size=14)
                rendered_using = "dynamic_content"
            clear_slide_prompt_text(slide)

        elif slide_type in ("agenda", "bullet", "summary", "closing"):
            items = slide_spec.get("items") or slide_spec.get(
                "bullet_points") or slide_spec.get("points") or []
            if not isinstance(items, list):
                items = [str(items)]
            used_placeholder = placeholder_policy != "ignore_placeholders" and add_bullet_content(
                slide, items)
            if not used_placeholder:
                body_text = "\n".join(f"- {item}" for item in items)
                ppt_utils.add_textbox(
                    slide, 0.9, 1.8, 8.0, 4.5, body_text, font_size=18)
                warnings.append(
                    "Bullet content used a dynamic textbox fallback")
                rendered_using = "dynamic_content"
            clear_slide_prompt_text(slide)

        elif slide_type == "two_column":
            left_title = slide_spec.get("left_title", "Left")
            right_title = slide_spec.get("right_title", "Right")
            left_points = slide_spec.get(
                "left_points") or slide_spec.get("left_items") or []
            right_points = slide_spec.get(
                "right_points") or slide_spec.get("right_items") or []

            left_text = left_title + "\n" + \
                "\n".join(f"- {item}" for item in left_points)
            right_text = right_title + "\n" + \
                "\n".join(f"- {item}" for item in right_points)

            ppt_utils.add_textbox(slide, 0.7, 1.7, 4.2,
                                  4.5, left_text, font_size=16)
            ppt_utils.add_textbox(slide, 5.0, 1.7, 4.2,
                                  4.5, right_text, font_size=16)
            rendered_using = "mixed"

        elif slide_type == "image_text":
            body_text = slide_spec.get(
                "content") or slide_spec.get("text") or ""
            image_path = slide_spec.get("image_path", "")
            body_placeholders = find_body_placeholders(slide)
            body_placeholder = body_placeholders[0] if body_placeholders else None
            if body_placeholder is not None and body_text:
                clear_text_frame(body_placeholder.text_frame)
                body_placeholder.text = str(body_text)
            else:
                ppt_utils.add_textbox(
                    slide, 0.7, 1.8, 4.2, 4.4, str(body_text), font_size=16)
                rendered_using = "dynamic_content"

            if image_path and os.path.exists(image_path):
                ppt_utils.add_image(slide, image_path, 5.2, 1.7, 4.0, 3.8)
            else:
                ppt_utils.add_textbox(slide, 5.2, 2.7, 3.5, 1.0, slide_spec.get(
                    "image_caption", "Image placeholder"), font_size=14)
                warnings.append(
                    "Image slide had no valid image path; inserted placeholder text")
                rendered_using = "mixed"

        elif slide_type == "table":
            table_data = slide_spec.get("table_data") or []
            if table_data and isinstance(table_data[0], list):
                rows = len(table_data)
                cols = max(len(row) for row in table_data)
                table_shape = ppt_utils.add_table(
                    slide, rows, cols, 0.7, 1.8, 8.5, 3.8)
                table = table_shape.table
                for row_index, row in enumerate(table_data):
                    for col_index, value in enumerate(row):
                        table.cell(row_index, col_index).text = str(value)
            else:
                warnings.append("Table slide received no valid table_data")

        elif slide_type == "chart":
            chart_data = slide_spec.get("chart_data") or {}
            categories = chart_data.get("categories") or ["Q1", "Q2", "Q3"]
            series = chart_data.get("series") or [
                {"name": "Series 1", "values": [1, 2, 3]}]
            series_names = [
                item.get("name", f"Series {idx + 1}") for idx, item in enumerate(series)]
            series_values = [item.get("values", []) for item in series]
            ppt_utils.add_chart(
                slide,
                chart_data.get("chart_type", "column"),
                0.8,
                1.8,
                8.2,
                3.8,
                categories,
                series_names,
                series_values,
            )
            rendered_using = "mixed"

        elif slide_type == "timeline":
            milestones = slide_spec.get(
                "milestones") or slide_spec.get("items") or []
            y = 2.1
            for index, item in enumerate(milestones[:5]):
                label = item if isinstance(item, str) else item.get(
                    "label", f"Milestone {index + 1}")
                detail = "" if isinstance(
                    item, str) else item.get("detail", "")
                ppt_utils.add_textbox(
                    slide, 0.8, y, 8.0, 0.6, f"{index + 1}. {label}", font_size=16, bold=True)
                if detail:
                    ppt_utils.add_textbox(
                        slide, 1.2, y + 0.4, 7.4, 0.5, detail, font_size=12)
                y += 1.0
            rendered_using = "dynamic_content"

        else:
            body_text = slide_spec.get("content") or slide_spec.get(
                "text") or slide_spec.get("notes") or ""
            if body_text:
                body_placeholders = find_body_placeholders(slide)
                body_placeholder = body_placeholders[0] if body_placeholders else None
                if body_placeholder is not None and placeholder_policy != "ignore_placeholders":
                    clear_text_frame(body_placeholder.text_frame)
                    body_placeholder.text = str(body_text)
                else:
                    ppt_utils.add_textbox(
                        slide, 0.8, 1.8, 8.2, 4.0, str(body_text), font_size=16)
                    rendered_using = "dynamic_content"

        clear_slide_prompt_text(slide)
        return {
            "slide_index": slide_index,
            "slide_type": slide_type,
            "layout_index": layout_index,
            "layout_name": getattr(layout, "name", f"Layout {layout_index}"),
            "rendered_using": rendered_using,
            "reused_template_slide": reused_template_slide,
            "warnings": warnings,
        }

    def standard_outline(topic: str, sections: List[str], page_count: int, scenario: str) -> List[Dict[str, Any]]:
        outline: List[Dict[str, Any]] = [
            {
                "slide_no": 1,
                "slide_type": "cover",
                "title": topic,
                "subtitle": "Generated from template-aware workflow",
            },
            {
                "slide_no": 2,
                "slide_type": "agenda",
                "title": "目录",
                "items": sections or ["背景", "分析", "方案", "落地计划"],
            },
        ]

        body_types = ["section", "bullet", "two_column", "chart", "summary"]
        for index, section in enumerate(sections or []):
            outline.append({
                "slide_no": len(outline) + 1,
                "slide_type": body_types[index % len(body_types)],
                "title": section,
                "items": [f"{section} 要点 1", f"{section} 要点 2", f"{section} 要点 3"],
            })

        while len(outline) < max(3, page_count - 1):
            item_index = len(outline) - 1
            outline.append({
                "slide_no": len(outline) + 1,
                "slide_type": body_types[item_index % len(body_types)],
                "title": f"{topic} - 内容页 {item_index}",
                "items": ["关键点 1", "关键点 2", "关键点 3"],
            })

        closing_title = "总结与下一步" if scenario != "closing_only" else topic
        outline.append({
            "slide_no": len(outline) + 1,
            "slide_type": "closing",
            "title": closing_title,
            "items": ["总结要点", "下一步计划", "Q&A"],
        })

        return outline[:max(3, page_count)]

    def merge_content(outline: List[Dict[str, Any]], content_spec: Optional[Dict[str, Any]]) -> List[Dict[str, Any]]:
        if not content_spec:
            return outline

        slide_updates: Dict[int, Dict[str, Any]] = {}
        for slide in content_spec.get("slides", []):
            slide_no = slide.get("slide_no")
            if isinstance(slide_no, int):
                slide_updates[slide_no] = slide

        merged: List[Dict[str, Any]] = []
        for slide in outline:
            combined = dict(slide)
            slide_no = slide.get("slide_no")
            if slide_no in slide_updates:
                combined.update(slide_updates[slide_no])
            merged.append(combined)

        return merged

    def profile_file_path(profile_name: str) -> str:
        profile_id = slugify(profile_name)
        return os.path.join(profile_dir, f"{profile_id}.json")

    def shape_text(shape) -> str:
        if hasattr(shape, "text_frame") and shape.text_frame:
            return shape.text_frame.text or ""
        return ""

    def shape_font_summary(shape) -> Dict[str, Any]:
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            return {}
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                return {
                    "font_name": font.name,
                    "font_size": font.size.pt if font.size else None,
                    "bold": font.bold,
                    "italic": font.italic,
                }
        return {}

    def summarize_shape(shape, shape_index: int) -> Dict[str, Any]:
        is_placeholder = bool(getattr(shape, "is_placeholder", False))
        placeholder = None
        if is_placeholder:
            placeholder = {
                "idx": getattr(shape.placeholder_format, "idx", None),
                "type": str(getattr(shape.placeholder_format, "type", "")),
                "role": placeholder_role(shape),
            }

        return {
            "shape_index": shape_index,
            "name": getattr(shape, "name", ""),
            "shape_type": str(getattr(shape, "shape_type", "")),
            "is_placeholder": is_placeholder,
            "placeholder": placeholder,
            "has_text": hasattr(shape, "text_frame") and shape.text_frame is not None,
            "text": shape_text(shape),
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
            "left_in": round(shape.left / emu_per_inch, 3),
            "top_in": round(shape.top / emu_per_inch, 3),
            "width_in": round(shape.width / emu_per_inch, 3),
            "height_in": round(shape.height / emu_per_inch, 3),
            "font": shape_font_summary(shape),
        }

    def is_textual_shape_summary(shape: Dict[str, Any]) -> bool:
        return bool(shape.get("has_text")) and (
            bool((shape.get("text") or "").strip())
            or "TEXT" in (shape.get("shape_type") or "")
            or shape.get("is_placeholder")
        )

    def candidate_score(shape: Dict[str, Any], role: str, slide_width: int, slide_height: int) -> float:
        text = (shape.get("text") or "").strip()
        normalized = normalize_text(text)
        left = shape.get("left", 0)
        top = shape.get("top", 0)
        width = max(shape.get("width", 0), 1)
        height = max(shape.get("height", 0), 1)
        area = width * height
        slide_area = max(slide_width * slide_height, 1)
        center_x = left + width / 2
        center_y = top + height / 2
        placeholder_role_name = (shape.get("placeholder") or {}).get("role")
        font_size = (shape.get("font") or {}).get("font_size") or 0

        score = 0.0
        if role == "title":
            score += 2.0 if placeholder_role_name == "title" else 0.0
            score += 1.5 if font_size >= 24 else 0.0
            score += 1.0 if area / slide_area > 0.015 else 0.0
            score += 1.0 if center_y < slide_height * 0.55 else 0.0
            score += 0.5 if len(text) <= 40 else -0.5
            score += 1.0 if any(token in normalized for token in [
                                "标题", "title", "目录", "content"]) else 0.0

        elif role == "subtitle":
            score += 2.0 if placeholder_role_name == "subtitle" else 0.0
            score += 1.0 if 10 <= font_size <= 24 else 0.0
            score += 1.0 if slide_height * 0.25 <= center_y <= slide_height * 0.75 else 0.0
            score += 0.7 if 10 <= len(text) <= 80 else 0.0
            score += 0.8 if any(token in normalized for token in [
                                "subtitle", "minimal", "template", "阐述", "说明"]) else 0.0

        elif role == "body":
            score += 2.0 if placeholder_role_name == "body" else 0.0
            score += 1.5 if area / slide_area > 0.02 else 0.0
            score += 1.0 if center_y > slide_height * 0.25 else 0.0
            score += 1.0 if len(text) >= 10 else 0.0
            score += 1.0 if any(token in normalized for token in [
                                "添加文字", "阐述", "说明", "content"]) else 0.0

        elif role == "agenda_item":
            score += 1.2 if re.fullmatch(r"\d{1,2}", text.strip()) else 0.0
            score += 1.0 if area / slide_area < 0.025 else 0.0
            score += 1.0 if slide_height * 0.25 <= center_y <= slide_height * 0.8 else 0.0
            score += 0.5 if center_x < slide_width * 0.85 else 0.0

        elif role == "clear":
            score += 2.0 if any(token in normalized for token in [
                                "20xx", "http", "免费网", "汇报人"]) else 0.0
            score += 1.5 if normalized.startswith("part") else 0.0
            score += 1.6 if normalized in {"content",
                                           "addpagetitlecontent"} else 0.0
            score += 1.0 if is_prompt_text(text) else 0.0

        return round(score, 3)

    def role_candidates_for_slide(slide_summary: Dict[str, Any], slide_width: int, slide_height: int) -> Dict[str, Any]:
        shapes = [shape for shape in slide_summary["shapes"]
                  if is_textual_shape_summary(shape)]
        roles = ["title", "subtitle", "body", "agenda_item", "clear"]
        candidates: Dict[str, Any] = {}

        for role in roles:
            scored = []
            for shape in shapes:
                score = candidate_score(shape, role, slide_width, slide_height)
                if score > 0:
                    scored.append({
                        "shape_index": shape["shape_index"],
                        "score": score,
                        "text": shape.get("text", ""),
                        "name": shape.get("name", ""),
                        "position": {
                            "left_in": shape.get("left_in"),
                            "top_in": shape.get("top_in"),
                            "width_in": shape.get("width_in"),
                            "height_in": shape.get("height_in"),
                        },
                    })
            scored.sort(key=lambda item: (-item["score"], item["shape_index"]))
            candidates[role] = scored[:8]

        return candidates

    def infer_slide_role(slide_summary: Dict[str, Any], candidates: Dict[str, Any]) -> str:
        index = slide_summary["slide_index"]
        layout_name = slide_summary.get("layout_name", "")
        texts = " ".join((shape.get("text") or "")
                         for shape in slide_summary.get("shapes", []))
        normalized = normalize_text(layout_name + " " + texts)

        if index == 0 or "封面" in normalized or "标题幻灯片" in layout_name:
            if "谢谢" in normalized or "thank" in normalized:
                return "closing"
            return "cover"
        if layout_name == "Main":
            return "summary"
        if "目录" in normalized or "content" in normalized:
            return "agenda"
        if "part" in normalized or "章" in normalized:
            return "section"
        if len(candidates.get("agenda_item", [])) >= 3:
            return "summary"
        if "谢谢" in normalized or "thank" in normalized:
            return "closing"
        return f"slide_{index + 1}"

    def build_draft_profile(slides: List[Dict[str, Any]], slide_width: int, slide_height: int) -> Dict[str, Any]:
        draft: Dict[str, Any] = {"slides": {}}
        used_roles = set()

        for slide_summary in slides:
            candidates = slide_summary.get("role_candidates", {})
            role = infer_slide_role(slide_summary, candidates)
            if role in used_roles:
                role = f"{role}_{slide_summary['slide_index'] + 1}"
            used_roles.add(role)

            title_candidate = (candidates.get("title")
                               or candidates.get("body") or [{}])[0]
            subtitle_candidate = (candidates.get(
                "subtitle") or candidates.get("body") or [{}])[0]
            body_candidate = (candidates.get("body")
                              or candidates.get("subtitle") or [{}])[0]
            clear_shapes = [item["shape_index"] for item in candidates.get(
                "clear", []) if item.get("score", 0) >= 1.5]

            fields: Dict[str, Any] = {}
            if role in ("cover", "closing"):
                if "shape_index" in title_candidate:
                    fields["title"] = {
                        "shape_index": title_candidate["shape_index"], "max_chars": 28}
                if "shape_index" in subtitle_candidate and subtitle_candidate["shape_index"] != fields.get("title", {}).get("shape_index"):
                    fields["subtitle"] = {
                        "shape_index": subtitle_candidate["shape_index"], "max_chars": 56}
            elif role == "agenda":
                title_options = candidates.get("title") or []
                agenda_title = next(
                    (item for item in title_options if "目录" in item.get("text", "")), title_candidate)
                if "shape_index" in agenda_title:
                    fields["title"] = {
                        "shape_index": agenda_title["shape_index"], "default": "目录"}
                item_candidates = candidates.get(
                    "agenda_item") or candidates.get("body") or []
                fields["items"] = [
                    {"shape_index": item["shape_index"], "max_chars": 12}
                    for item in item_candidates[:4]
                ]
            elif role == "section":
                if "shape_index" in title_candidate:
                    fields["title"] = {
                        "shape_index": title_candidate["shape_index"], "max_chars": 18}
                if "shape_index" in body_candidate:
                    fields["body"] = {
                        "shape_index": body_candidate["shape_index"], "max_chars": 90}
            else:
                if "shape_index" in title_candidate:
                    fields["title"] = {
                        "shape_index": title_candidate["shape_index"], "max_chars": 18}
                item_candidates = candidates.get(
                    "body") or candidates.get("agenda_item") or []
                fields["items"] = [
                    {"shape_index": item["shape_index"], "max_chars": 18}
                    for item in item_candidates[:4]
                ]

            used_shape_indices = set()
            for field_spec in fields.values():
                if isinstance(field_spec, dict) and isinstance(field_spec.get("shape_index"), int):
                    used_shape_indices.add(field_spec["shape_index"])
                elif isinstance(field_spec, list):
                    for item in field_spec:
                        if isinstance(item, dict) and isinstance(item.get("shape_index"), int):
                            used_shape_indices.add(item["shape_index"])
            clear_shapes = [
                shape_index for shape_index in clear_shapes if shape_index not in used_shape_indices]

            draft["slides"][role] = {
                "source_slide_index": slide_summary["slide_index"],
                "clear_shapes": clear_shapes,
                "fields": fields,
            }

        return draft

    def resolve_profile_template_path(profile: Dict[str, Any]) -> Tuple[Optional[str], Optional[str]]:
        template_path = profile.get("template_path", "")
        if template_path and os.path.exists(template_path):
            return template_path, None

        template_name = profile.get("template_name", "")
        resolved_path, resolution = resolve_template_reference(
            {"template_name": template_name})
        if resolved_path:
            return resolved_path, None
        return None, resolution.get("error", f"Template not found: {template_name}")

    def remove_unprofiled_slides(presentation, keep_indices: List[int]) -> int:
        keep = set(keep_indices)
        removed = 0
        index = len(presentation.slides) - 1
        while index >= 0:
            if index not in keep:
                slide_id = presentation.slides._sldIdLst[index]
                relationship_id = slide_id.rId
                presentation.part.drop_rel(relationship_id)
                presentation.slides._sldIdLst.remove(slide_id)
                removed += 1
            index -= 1
        return removed

    def get_profile_content(content: Dict[str, Any], slide_role: str) -> Dict[str, Any]:
        slides = content.get("slides", {})
        if isinstance(slides, dict):
            slide_content = slides.get(slide_role, {})
            return slide_content if isinstance(slide_content, dict) else {}

        if isinstance(slides, list):
            for item in slides:
                if isinstance(item, dict) and item.get("slide_role") == slide_role:
                    return item
        return {}

    def get_field_value(content: Dict[str, Any], slide_content: Dict[str, Any], field_name: str, default: Any = "") -> Any:
        if field_name in slide_content:
            return slide_content[field_name]
        if field_name in content:
            return content[field_name]
        return default

    def has_field_value(content: Dict[str, Any], slide_content: Dict[str, Any], field_name: str) -> bool:
        if field_name in slide_content:
            value = slide_content[field_name]
        elif field_name in content:
            value = content[field_name]
        else:
            return False

        if isinstance(value, str):
            return bool(value.strip())
        if isinstance(value, list):
            return bool(value)
        return value is not None

    def missing_profile_fields(slide_mapping: Dict[str, Any], content: Dict[str, Any], slide_content: Dict[str, Any]) -> List[str]:
        missing = []
        for field_name, field_spec in slide_mapping.get("fields", {}).items():
            if isinstance(field_spec, dict) and "default" in field_spec:
                continue
            if not has_field_value(content, slide_content, field_name):
                missing.append(field_name)
        return missing

    def write_profile_field(slide, field_name: str, field_spec: Any, content: Dict[str, Any], slide_content: Dict[str, Any]) -> Dict[str, Any]:
        if isinstance(field_spec, int):
            value = get_field_value(content, slide_content, field_name)
            if 0 <= field_spec < len(slide.shapes):
                set_shape_text(slide.shapes[field_spec], str(value))
                return {"field": field_name, "shape_index": field_spec, "value": value}
            return {"field": field_name, "shape_index": field_spec, "error": "shape_index out of range"}

        if isinstance(field_spec, dict):
            shape_index = field_spec.get("shape_index")
            value = get_field_value(
                content, slide_content, field_name, field_spec.get("default", ""))
            max_chars = field_spec.get("max_chars")
            if isinstance(value, str) and isinstance(max_chars, int) and max_chars > 0:
                value = value[:max_chars]
            if isinstance(shape_index, int) and 0 <= shape_index < len(slide.shapes):
                set_shape_text(slide.shapes[shape_index], str(value))
                apply_field_format(slide.shapes[shape_index], field_spec)
                return {"field": field_name, "shape_index": shape_index, "value": value}
            return {"field": field_name, "shape_index": shape_index, "error": "shape_index out of range"}

        if isinstance(field_spec, list):
            values = get_field_value(content, slide_content, field_name, [])
            if not isinstance(values, list):
                values = [values]

            writes = []
            for index, item_spec in enumerate(field_spec):
                if isinstance(item_spec, int):
                    shape_index = item_spec
                    max_chars = None
                elif isinstance(item_spec, dict):
                    shape_index = item_spec.get("shape_index")
                    max_chars = item_spec.get("max_chars")
                else:
                    writes.append(
                        {"field": field_name, "item_index": index, "error": "invalid field item spec"})
                    continue

                value = values[index] if index < len(values) else ""
                if isinstance(value, str) and isinstance(max_chars, int) and max_chars > 0:
                    value = value[:max_chars]
                if isinstance(shape_index, int) and 0 <= shape_index < len(slide.shapes):
                    set_shape_text(slide.shapes[shape_index], str(value))
                    if isinstance(item_spec, dict):
                        apply_field_format(
                            slide.shapes[shape_index], item_spec)
                    writes.append({"field": field_name, "item_index": index,
                                  "shape_index": shape_index, "value": value})
                else:
                    writes.append({"field": field_name, "item_index": index,
                                  "shape_index": shape_index, "error": "shape_index out of range"})
            return {"field": field_name, "items": writes}

        return {"field": field_name, "error": "unsupported field spec"}

    @app.custom_route("/downloads/{filename}", methods=["GET"])
    async def download_presentation(request: Request):
        """Serve exported presentations for simplified workflow mode."""
        filename = os.path.basename(request.path_params["filename"])
        file_path = os.path.join(download_dir, filename)

        if not os.path.exists(file_path):
            for project in projects.values():
                if not isinstance(project, dict):
                    continue
                exported_path = project.get("last_export_path", "")
                if exported_path and os.path.basename(exported_path) == filename and os.path.exists(exported_path):
                    file_path = exported_path
                    break
            else:
                return JSONResponse({"error": f"File not found: {filename}"}, status_code=404)

        return FileResponse(
            path=file_path,
            filename=filename,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    @app.tool(
        annotations=ToolAnnotations(
            title="List Presentation Options",
            readOnlyHint=True,
        ),
    )
    def list_presentation_options() -> Dict[str, Any]:
        """List available visual themes and slide layout schemas for direct presentation generation."""
        return {
            "themes": [
                {
                    "theme_id": theme_id,
                    "name": theme["name"],
                    "description": theme["description"],
                    "font_name": theme["font_name"],
                    "colors": theme["colors"],
                }
                for theme_id, theme in get_builtin_themes().items()
            ],
            "default_theme": "business_blue",
            "layouts": get_builtin_layouts(),
            "default_layout_sequence": [
                "cover",
                "summary",
                "research_questions",
                "literature_matrix",
                "theoretical_framework",
                "method_design",
                "findings",
                "contribution_limitations",
                "closing",
            ],
            "academic_recommended_theme": "academic_burgundy",
            "common_slide_fields": ["type", "title", "points", "source_note"],
            "compatible_content_fields": ["evidence", "explanation", "analysis", "result", "conclusion", "mechanism", "boundary"],
            "compatible_source_fields": ["source_refs", "source", "source_text", "citation", "reference"],
            "auto_split_rules": {
                "table": "Split table rows into multiple slides when rows exceed 6.",
                "literature_matrix": "Split literature/studies/items into multiple slides when rows exceed 6.",
                "cards": "Split items/sections into multiple slides when item count exceeds 6.",
                "findings": "Split findings into multiple slides when finding count exceeds 4, and split long finding points into continuation cards.",
                "contribution_limitations": "Split contributions, limitations, and future directions into multiple slides when each group exceeds 3 items.",
            },
            "usage_hint": "Choose one theme for the deck, then set slides[].type to a layout_id. Prefer simple slide fields: type, title, points, and source_note. The generator also accepts evidence/explanation/source_refs and normalizes them internally.",
        }

    @app.tool(
        annotations=ToolAnnotations(
            title="Generate Presentation",
        ),
    )
    def generate_presentation(
        presentation_id: str,
        title: str,
        subtitle: str = "",
        slides: List[Dict[str, Any]] = [],
        theme: str = "business_blue",
        density: str = "standard",
        overflow: str = "shrink_then_truncate",
        style: str = "business",
        visual_level: str = "clean",
        language: str = "zh-CN",
        page_size: str = "wide_16_9",
        auto_cover: bool = True,
        auto_closing: bool = False,
        show_page_number: bool = True,
        show_footer: bool = True,
        footer_text: str = "",
        output_name: str = "",
    ) -> Dict[str, Any]:
        """Generate a complete PPT directly with python-pptx, without using a template.

        Supported slide types: cover, summary, cards, comparison, process,
        timeline, metrics, architecture, table, quote, and closing.
        If a slide omits type, the tool infers it from fields such as
        items, steps, metrics, layers, table, comparisons, or statement.
        """
        if not presentation_id:
            presentation_id = f"deck_{slugify(title)}"

        effective_theme_id, theme_spec = get_theme(theme)
        presentation = ppt_utils.create_presentation()
        effective_page_size = apply_page_size(presentation, page_size)

        if hasattr(ppt_utils, "set_core_properties"):
            ppt_utils.set_core_properties(presentation, title=title)

        normalized_style = (style or "business").strip().lower()
        normalized_visual_level = (visual_level or "clean").strip().lower()
        if normalized_visual_level not in {"clean", "rich", "dense"}:
            normalized_visual_level = "clean"

        slide_specs = prepare_slide_specs(
            slides,
            title,
            subtitle,
            auto_cover,
            auto_closing,
            normalized_style,
        )
        slide_specs = expand_capacity_slide_specs(slide_specs)

        rendered_slide_types = []
        ignored_fields = []
        warnings: List[str] = []
        for index, slide_spec in enumerate(slide_specs):
            if not isinstance(slide_spec, dict):
                slide_spec = {"type": "summary", "title": str(slide_spec)}
            slide_spec = normalize_slide_spec(slide_spec, title, subtitle)
            if index == 0:
                slide_spec.setdefault("title", title)
                slide_spec.setdefault("subtitle", subtitle)
            slide_density = (slide_spec.get("density")
                             or density or "standard").strip().lower()
            slide_overflow = (slide_spec.get(
                "overflow") or overflow or "shrink_then_truncate").strip().lower()
            rendered_type = render_generated_slide(
                presentation,
                slide_spec,
                theme_spec,
                slide_density,
                slide_overflow,
                warnings,
            )
            rendered_slide_types.append(rendered_type)
            ignored_fields.extend(collect_ignored_fields(
                slide_spec, rendered_type, index + 1))

        effective_footer_text = footer_text or title
        add_deck_footer(
            presentation,
            theme_spec,
            effective_footer_text,
            show_footer,
            show_page_number,
            normalized_visual_level,
        )
        quality = inspect_presentation_quality(presentation, warnings)

        presentations[presentation_id] = presentation
        projects[presentation_id] = {
            "title": title,
            "subtitle": subtitle,
            "output_name": output_name or f"{slugify(title)}.pptx",
            "generation_mode": "direct_python_pptx",
            "theme": effective_theme_id,
            "density": density,
            "overflow": overflow,
            "style": normalized_style,
            "visual_level": normalized_visual_level,
            "language": language,
            "page_size": effective_page_size,
        }
        set_current_presentation_id(presentation_id)

        return {
            "presentation_id": presentation_id,
            "title": title,
            "slide_count": len(presentation.slides),
            "theme": effective_theme_id,
            "style": normalized_style,
            "visual_level": normalized_visual_level,
            "page_size": effective_page_size,
            "rendered_slide_types": rendered_slide_types,
            "ignored_fields": ignored_fields,
            "warnings": warnings,
            "quality": quality,
            "next_step": "Call export_presentation with this presentation_id to save and get a download_url.",
        }

    @app.tool(
        annotations=ToolAnnotations(
            title="List Templates",
            readOnlyHint=True,
        ),
    )
    def list_templates(
        template_query: str = "",
        template_directory: str = "",
        include_layouts: bool = True,
        include_placeholders: bool = False,
    ) -> Dict[str, Any]:
        """List available PowerPoint templates from the configured template directories."""
        template_files = collect_template_files(template_directory)
        query = normalize_text(template_query)
        templates = []

        for template_path in template_files:
            template_summary = summarize_template_file(
                template_path, include_layouts, include_placeholders)
            haystack = normalize_text(
                " ".join([
                    template_summary["template_name"],
                    " ".join(template_summary.get("style_tags", [])),
                    template_summary.get(
                        "core_properties", {}).get("title") or "",
                    template_summary.get(
                        "core_properties", {}).get("subject") or "",
                ])
            )
            if query and query not in haystack:
                continue
            templates.append(template_summary)

        return {
            "templates": templates,
            "total_templates": len(templates),
            "searched_directories": sorted(set(os.path.dirname(path) for path in template_files)),
        }

    @app.tool(
        annotations=ToolAnnotations(
            title="Analyze Template",
            readOnlyHint=True,
        ),
    )
    def analyze_template(
        template_path: str = "",
        template_name: str = "",
        template_directory: str = "",
        include_non_text_shapes: bool = False,
        include_role_candidates: bool = True,
        include_draft_profile: bool = True,
        max_slides: int = 20,
    ) -> Dict[str, Any]:
        """Analyze a PowerPoint template and return slide/shape metadata for profile mapping."""
        resolved_path, resolution = resolve_template_reference({
            "template_path": template_path,
            "template_name": template_name,
            "template_directory": template_directory,
        })
        if not resolved_path:
            return {"error": resolution.get("error", "Template not found")}

        try:
            presentation = ppt_utils.open_presentation(resolved_path)
            slides = []
            for slide_index, slide in enumerate(presentation.slides):
                if slide_index >= max_slides:
                    break

                shapes = []
                for shape_index, shape in enumerate(slide.shapes):
                    summary = summarize_shape(shape, shape_index)
                    if not include_non_text_shapes and not summary["has_text"] and not summary["is_placeholder"]:
                        continue
                    shapes.append(summary)

                slides.append({
                    "slide_index": slide_index,
                    "layout_name": getattr(slide.slide_layout, "name", ""),
                    "shape_count": len(slide.shapes),
                    "shapes": shapes,
                })

            if include_role_candidates:
                for slide_summary in slides:
                    slide_summary["role_candidates"] = role_candidates_for_slide(
                        slide_summary,
                        presentation.slide_width,
                        presentation.slide_height,
                    )
                    slide_summary["suggested_slide_role"] = infer_slide_role(
                        slide_summary,
                        slide_summary["role_candidates"],
                    )

            draft_profile = None
            if include_draft_profile:
                draft_profile = build_draft_profile(
                    slides,
                    presentation.slide_width,
                    presentation.slide_height,
                )

            return {
                "template_name": os.path.basename(resolved_path),
                "template_path": resolved_path,
                "slide_count": len(presentation.slides),
                "layout_count": len(presentation.slide_layouts),
                "slides": slides,
                "draft_profile": draft_profile,
                "profile_hint": {
                    "profile_name": slugify(os.path.splitext(os.path.basename(resolved_path))[0]),
                    "next_step": "Review draft_profile, adjust field shape_index values if needed, then pass it as mapping to create_template_profile."
                },
            }
        except Exception as exc:
            return {"error": f"Failed to analyze template: {str(exc)}"}

    @app.tool(
        annotations=ToolAnnotations(
            title="Create Template Profile",
        ),
    )
    def create_template_profile(
        profile_name: str,
        template_path: str = "",
        template_name: str = "",
        mapping: Optional[Dict[str, Any]] = None,
        description: str = "",
        overwrite: bool = False,
    ) -> Dict[str, Any]:
        """Create or update a reusable template profile mapping slide roles to shape fields."""
        if not profile_name:
            return {"error": "profile_name is required"}

        resolved_path, resolution = resolve_template_reference({
            "template_path": template_path,
            "template_name": template_name,
        })
        if not resolved_path:
            return {"error": resolution.get("error", "Template not found")}

        mapping = mapping or {}
        if "slides" not in mapping or not isinstance(mapping["slides"], dict):
            return {"error": "mapping.slides must be a dictionary of slide roles"}

        path = profile_file_path(profile_name)
        if os.path.exists(path) and not overwrite:
            return {
                "error": f"Template profile already exists: {profile_name}. Pass overwrite=true to replace it.",
                "profile_path": path,
            }

        profile = {
            "profile_name": profile_name,
            "profile_id": slugify(profile_name),
            "description": description,
            "template_name": os.path.basename(resolved_path),
            "template_path": resolved_path,
            "mapping": mapping,
        }

        with open(path, "w", encoding="utf-8") as file:
            json.dump(profile, file, ensure_ascii=False, indent=2)

        return {
            "message": "Template profile saved",
            "profile_name": profile_name,
            "profile_id": profile["profile_id"],
            "profile_path": path,
            "template_path": resolved_path,
            "slide_roles": list(mapping["slides"].keys()),
        }

    @app.tool(
        annotations=ToolAnnotations(
            title="Generate From Template Profile",
        ),
    )
    def generate_from_template_profile(
        profile_name: str,
        content: Dict[str, Any],
        presentation_id: str = "",
        output_name: str = "",
        keep_unmapped_slides: bool = False,
        slide_roles: Optional[List[str]] = None,
        require_content: bool = True,
    ) -> Dict[str, Any]:
        """Generate a presentation by filling a saved template profile."""
        path = profile_file_path(profile_name)
        if not os.path.exists(path):
            return {"error": f"Template profile not found: {profile_name}", "profile_path": path}

        with open(path, "r", encoding="utf-8") as file:
            profile = json.load(file)

        template_path, template_error = resolve_profile_template_path(profile)
        if template_error:
            return {"error": template_error}

        try:
            presentation = ppt_utils.create_presentation_from_template(
                template_path)
            mapping = profile.get("mapping", {})
            slide_mappings = mapping.get("slides", {})
            if slide_roles:
                missing_roles = [
                    role for role in slide_roles if role not in slide_mappings]
                if missing_roles:
                    return {
                        "error": f"Template profile does not define requested slide roles: {missing_roles}",
                        "available_slide_roles": list(slide_mappings.keys()),
                    }
                selected_slide_mappings = {
                    role: slide_mappings[role]
                    for role in slide_roles
                }
            else:
                selected_slide_mappings = slide_mappings

            used_slide_indices: List[int] = []
            field_results = []

            missing_content = []
            if require_content:
                for slide_role, slide_mapping in selected_slide_mappings.items():
                    slide_content = get_profile_content(content, slide_role)
                    missing_fields = missing_profile_fields(
                        slide_mapping, content, slide_content)
                    if missing_fields:
                        missing_content.append({
                            "slide_role": slide_role,
                            "missing_fields": missing_fields,
                        })

            if missing_content:
                return {
                    "error": "Missing required content for template profile fields.",
                    "missing_content": missing_content,
                    "hint": "Pass slide_roles to generate only the roles you want, or provide content for each selected role.",
                }

            for slide_role, slide_mapping in selected_slide_mappings.items():
                source_slide_index = slide_mapping.get("source_slide_index")
                if not isinstance(source_slide_index, int) or source_slide_index < 0 or source_slide_index >= len(presentation.slides):
                    field_results.append({
                        "slide_role": slide_role,
                        "error": f"Invalid source_slide_index: {source_slide_index}",
                    })
                    continue

                used_slide_indices.append(source_slide_index)
                slide = presentation.slides[source_slide_index]
                slide_content = get_profile_content(content, slide_role)
                cleared_shapes = []
                for shape_index in slide_mapping.get("clear_shapes", []):
                    if isinstance(shape_index, int) and 0 <= shape_index < len(slide.shapes):
                        set_shape_text(slide.shapes[shape_index], "")
                        cleared_shapes.append(shape_index)

                fields = slide_mapping.get("fields", {})
                writes = []
                for field_name, field_spec in fields.items():
                    writes.append(write_profile_field(
                        slide, field_name, field_spec, content, slide_content))

                if slide_mapping.get("clear_prompt_text", True):
                    clear_slide_prompt_text(slide)

                field_results.append({
                    "slide_role": slide_role,
                    "source_slide_index": source_slide_index,
                    "cleared_shapes": cleared_shapes,
                    "writes": writes,
                })

            removed_slides = 0
            if not keep_unmapped_slides:
                removed_slides = remove_unprofiled_slides(
                    presentation, used_slide_indices)

            effective_presentation_id = presentation_id or f"profile_{slugify(profile_name)}"
            effective_output_name = output_name or content.get(
                "output_name") or f"{effective_presentation_id}.pptx"
            presentations[effective_presentation_id] = presentation
            projects[effective_presentation_id] = {
                "title": content.get("title", profile_name),
                "output_name": effective_output_name,
                "template_profile": profile_name,
                "resolved_template_path": template_path,
                "profile_path": path,
            }
            set_current_presentation_id(effective_presentation_id)

            return {
                "presentation_id": effective_presentation_id,
                "message": "Presentation generated from template profile",
                "profile_name": profile_name,
                "template_path": template_path,
                "slide_count": len(presentation.slides),
                "removed_unmapped_slides": removed_slides,
                "field_results": field_results,
            }
        except Exception as exc:
            return {"error": f"Failed to generate from template profile: {str(exc)}"}

    @app.tool(
        annotations=ToolAnnotations(
            title="Create Presentation Project",
        ),
    )
    def create_presentation_project(
        title: str,
        subtitle: str = "",
        language: str = "zh-CN",
        output_name: str = "",
        template: Optional[Dict[str, Any]] = None,
        presentation_id: str = "",
    ) -> Dict[str, Any]:
        """Create a new template-aware presentation project."""
        template = template or {}
        presentation_id = presentation_id or f"proj_{slugify(title)}"
        resolved_template_path, resolution = resolve_template_reference(
            template)

        if resolution.get("error"):
            return resolution

        try:
            if resolved_template_path:
                presentation = ppt_utils.create_presentation_from_template(
                    resolved_template_path)
            else:
                presentation = ppt_utils.create_presentation()

            ppt_utils.set_core_properties(
                presentation,
                title=title,
                subject=subtitle or title,
                comments=f"language={language}",
            )

            presentations[presentation_id] = presentation
            projects[presentation_id] = {
                "title": title,
                "subtitle": subtitle,
                "language": language,
                "output_name": output_name or f"{slugify(title)}.pptx",
                "template": dict(template),
                "resolved_template_path": resolved_template_path,
                "template_resolution": resolution,
                "outline": [],
                "template_slide_count": len(presentation.slides),
            }
            set_current_presentation_id(presentation_id)

            resolved_template = None
            if resolved_template_path:
                resolved_template = {
                    "template_name": os.path.basename(resolved_template_path),
                    "template_path": resolved_template_path,
                    "mode": resolution.get("mode", "strong_template"),
                    "selection_method": resolution.get("selection_method"),
                }

            return {
                "presentation_id": presentation_id,
                "message": "Presentation project created",
                "resolved_template": resolved_template,
                "layout_catalog": build_layout_catalog(presentation),
                "project": projects[presentation_id],
            }
        except Exception as exc:
            return {"error": f"Failed to create presentation project: {str(exc)}"}

    @app.tool(
        annotations=ToolAnnotations(
            title="Plan Presentation",
        ),
    )
    def plan_presentation(
        presentation_id: str,
        topic: str,
        audience: str = "general",
        page_count: int = 8,
        scenario: str = "proposal",
        template_context: Optional[Dict[str, Any]] = None,
        sections: Optional[List[str]] = None,
    ) -> Dict[str, Any]:
        """Generate a template-aware page outline for a presentation project."""
        if presentation_id not in presentations:
            return {"error": "No presentation is currently loaded or the specified ID is invalid"}

        template_context = template_context or {}
        sections = sections or []
        presentation = presentations[presentation_id]
        outline = standard_outline(topic, sections, page_count, scenario)

        if template_context.get("use_template_layouts", True):
            for slide in outline:
                layout_index = select_layout_index(
                    presentation,
                    slide["slide_type"],
                    None,
                    "preferred_then_best_match",
                )
                layout = presentation.slide_layouts[layout_index]
                slide["preferred_layout"] = summarize_layout(
                    layout,
                    layout_index,
                    include_placeholders=template_context.get(
                        "prefer_detected_layout_mapping", True),
                )
                slide["placeholder_map"] = layout_placeholder_map(layout)

        projects.setdefault(presentation_id, {})
        projects[presentation_id]["outline"] = outline
        projects[presentation_id]["planning_context"] = {
            "topic": topic,
            "audience": audience,
            "page_count": page_count,
            "scenario": scenario,
        }

        return {
            "presentation_id": presentation_id,
            "audience": audience,
            "outline": outline,
        }

    @app.tool(
        annotations=ToolAnnotations(
            title="Build Presentation",
        ),
    )
    def build_presentation(
        presentation_id: str,
        outline: Optional[List[Dict[str, Any]]] = None,
        content_spec: Optional[Dict[str, Any]] = None,
        template_rendering: Optional[Dict[str, Any]] = None,
        reset_existing_slides: bool = True,
    ) -> Dict[str, Any]:
        """Build a complete PowerPoint deck from a high-level outline and content spec."""
        if presentation_id not in presentations:
            return {"error": "No presentation is currently loaded or the specified ID is invalid"}

        presentation = presentations[presentation_id]
        template_rendering = template_rendering or {}
        outline = outline or projects.get(
            presentation_id, {}).get("outline") or []

        if not outline:
            return {"error": "No outline provided. Use plan_presentation first or pass outline directly."}

        require_planned_outline = template_rendering.get(
            "require_planned_outline", True)
        if require_planned_outline and any("preferred_layout" not in slide for slide in outline):
            return {
                "error": "Build requires a planned outline with preferred_layout metadata. Call plan_presentation first, then pass its outline into build_presentation."
            }

        merged_outline = merge_content(outline, content_spec)

        removed_slides = 0
        project = projects.get(presentation_id, {})
        existing_template_slide_count = project.get("template_slide_count", 0)
        reuse_existing_template_slides = template_rendering.get(
            "use_template_sample_slides", True) and existing_template_slide_count > 0
        if reset_existing_slides and not reuse_existing_template_slides:
            removed_slides = clear_all_slides(presentation)

        template_rendering = dict(template_rendering)
        template_rendering["existing_template_slide_count"] = existing_template_slide_count
        build_results = []
        warnings: List[str] = []
        for slide_position, slide_spec in enumerate(merged_outline):
            result = render_slide_from_spec(
                presentation, slide_spec, template_rendering, slide_position)
            build_results.append(result)
            warnings.extend(result["warnings"])

        removed_template_tail_slides = 0
        if reuse_existing_template_slides and len(presentation.slides) > len(merged_outline):
            removed_template_tail_slides = remove_slides_from(
                presentation, len(merged_outline))

        projects.setdefault(presentation_id, {})
        projects[presentation_id]["outline"] = merged_outline

        strong_template_slides = len(
            [item for item in build_results if item["rendered_using"] == "template_layout"])
        mixed_slides = len(
            [item for item in build_results if item["rendered_using"] == "mixed"])
        dynamic_slides = len(
            [item for item in build_results if item["rendered_using"] == "dynamic_content"])

        return {
            "presentation_id": presentation_id,
            "slide_count": len(presentation.slides),
            "removed_existing_slides": removed_slides,
            "removed_template_tail_slides": removed_template_tail_slides,
            "build_results": build_results,
            "build_summary": {
                "strong_template_slides": strong_template_slides,
                "mixed_slides": mixed_slides,
                "dynamic_slides": dynamic_slides,
                "fallback_applied": dynamic_slides > 0 or mixed_slides > 0,
            },
            "warnings": warnings,
        }

    @app.tool(
        annotations=ToolAnnotations(
            title="Revise Presentation",
        ),
    )
    def revise_presentation(
        presentation_id: str,
        instructions: Optional[List[str]] = None,
        changes: Optional[List[Dict[str, Any]]] = None,
        template_revision: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        """Apply a limited set of high-level revisions to an existing presentation."""
        if presentation_id not in presentations:
            return {"error": "No presentation is currently loaded or the specified ID is invalid"}

        presentation = presentations[presentation_id]
        instructions = instructions or []
        changes = changes or []
        template_revision = template_revision or {}

        applied_changes: List[Dict[str, Any]] = []
        warnings: List[str] = []

        for instruction in instructions:
            matched = re.match(r"第\s*(\d+)\s*页标题改为[:：]?\s*(.+)", instruction)
            if matched:
                slide_no = int(matched.group(1))
                if 1 <= slide_no <= len(presentation.slides):
                    slide = presentation.slides[slide_no - 1]
                    fill_title_placeholders(
                        slide, title=matched.group(2).strip())
                    applied_changes.append(
                        {"type": "rename_title", "slide_no": slide_no})
                    continue
                warnings.append(
                    f"Instruction skipped because slide {slide_no} does not exist: {instruction}")
                continue

            matched = re.match(r"删除第\s*(\d+)\s*页", instruction)
            if matched:
                slide_no = int(matched.group(1))
                if 1 <= slide_no <= len(presentation.slides):
                    slide_id = presentation.slides._sldIdLst[slide_no - 1]
                    relationship_id = slide_id.rId
                    presentation.part.drop_rel(relationship_id)
                    presentation.slides._sldIdLst.remove(slide_id)
                    applied_changes.append(
                        {"type": "delete_slide", "slide_no": slide_no})
                    continue
                warnings.append(
                    f"Instruction skipped because slide {slide_no} does not exist: {instruction}")
                continue

            matched = re.match(r"整体字体改为[:：]?\s*(.+)", instruction)
            if matched:
                font_name = matched.group(1).strip()
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            ppt_utils.format_text_advanced(
                                shape.text_frame, font_name=font_name)
                applied_changes.append(
                    {"type": "global_font", "font_name": font_name})
                continue

            warnings.append(
                f"Instruction not recognized and was skipped: {instruction}")

        for change in changes:
            slide_no = change.get("slide_no")
            if not isinstance(slide_no, int) or not (1 <= slide_no <= len(presentation.slides)):
                warnings.append(
                    f"Structured change skipped due to invalid slide_no: {change}")
                continue

            slide = presentation.slides[slide_no - 1]
            if change.get("action") == "set_title":
                fill_title_placeholders(slide, title=change.get("title", ""))
                applied_changes.append(
                    {"type": "set_title", "slide_no": slide_no})
            elif change.get("action") == "set_notes":
                notes = change.get("text", "")
                ppt_utils.add_textbox(
                    slide, 0.8, 6.6, 8.0, 0.4, notes, font_size=10)
                applied_changes.append(
                    {"type": "set_notes", "slide_no": slide_no})
            else:
                warnings.append(
                    f"Structured change action is not supported yet: {change.get('action')}")

        return {
            "presentation_id": presentation_id,
            "revised_slides": sorted({item.get("slide_no") for item in applied_changes if item.get("slide_no")}),
            "style_preserved": template_revision.get("preserve_template_style", True),
            "applied_changes": applied_changes,
            "warnings": warnings,
        }

    @app.tool(
        annotations=ToolAnnotations(
            title="Export Presentation",
        ),
    )
    def export_presentation(
        presentation_id: str,
        file_name: str = "",
        output_directory: str = "",
        template_export: Dict[str, Any] = {},
    ) -> Dict[str, Any]:
        """Save and export a presentation built through the simplified workflow."""
        if presentation_id not in presentations:
            return {"error": "No presentation is currently loaded or the specified ID is invalid"}

        template_export = template_export or {}
        project = projects.get(presentation_id, {})
        export_directory = output_directory or os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
            "ppt",
        )
        os.makedirs(export_directory, exist_ok=True)

        effective_file_name = file_name or project.get(
            "output_name") or f"{presentation_id}.pptx"
        if not effective_file_name.lower().endswith(".pptx"):
            effective_file_name = f"{effective_file_name}.pptx"

        file_path = os.path.join(export_directory, effective_file_name)
        try:
            ppt_utils.save_presentation(
                presentations[presentation_id], file_path)
        except PermissionError:
            return {
                "error": f"Cannot save presentation because the target file is locked or not writable: {file_path}",
                "hint": "Close the file if it is open, or export with a different file_name.",
                "presentation_id": presentation_id,
                "file_path": file_path,
            }
        project["last_export_path"] = file_path

        port = getattr(app.settings, "port", 8000)
        return {
            "presentation_id": presentation_id,
            "file_path": file_path,
            "download_url": f"{(download_url or f'http://localhost:{port}').rstrip('/')}/downloads/{effective_file_name}",
            "preserve_theme_assets": template_export.get("preserve_theme_assets", True),
        }

    template_profile_tools_enabled = (
        is_env_enabled("PPT_ENABLE_TEMPLATE_PROFILE_TOOLS")
        or is_env_enabled("PPT_ENABLE_ADVANCED_TOOLS")
    )
    if not template_profile_tools_enabled:
        for tool_name in [
            "list_templates",
            "analyze_template",
            "create_template_profile",
            "generate_from_template_profile",
        ]:
            app.remove_tool(tool_name)

    legacy_workflow_enabled = (
        is_env_enabled("PPT_ENABLE_LEGACY_WORKFLOW_TOOLS")
        or is_env_enabled("PPT_ENABLE_ADVANCED_TOOLS")
    )
    if not legacy_workflow_enabled:
        for tool_name in [
            "create_presentation_project",
            "plan_presentation",
            "build_presentation",
            "revise_presentation",
        ]:
            app.remove_tool(tool_name)
